"""
PPTX 시각적 겹침 탐지 v2 — 박스 좌표 + 폰트 실제 높이 기반.

v1 한계: 박스 좌표만 비교 → 큰 폰트가 박스 밖으로 튀어나오는 경우 못 잡음.
v2 개선:
  1. 박스 좌표 비교
  2. 텍스트 박스의 실제 폰트 크기를 XML에서 읽어 글자 높이(인치) 계산
  3. vertical_anchor가 'top'이면 글자가 박스 위에서 시작 → 실제 글자 영역 = (y, y + font_h * line_count)
     'middle'이면 박스 중앙 정렬 → 실제 글자 영역 = (y + (h - font_h*lc)/2, ...)
  4. "박스 좌표"와 "글자 실제 영역" 둘 다로 겹침 탐지
"""

import sys
import re
from pathlib import Path
from pptx import Presentation
from pptx.enum.text import MSO_ANCHOR

EMU_PER_INCH = 914400
PT_PER_INCH = 72

def emu_to_inch(v):
    return v / EMU_PER_INCH

def font_pt_to_inch(pt):
    """폰트 1줄 실제 높이(어센더+디센더 포함) ≈ pt × 1.0 / 72.
    Pretendard Black은 디센더 좀 길어 1.1배 안전 마진 적용."""
    return pt / PT_PER_INCH * 1.1

def shape_text(shape):
    if not shape.has_text_frame:
        return ""
    txt = shape.text_frame.text or ""
    txt = txt.replace("\n", " ").strip()
    return txt[:35] + ("…" if len(txt) > 35 else "")

def get_max_font_size(shape):
    """텍스트 박스 안 모든 run의 최대 폰트 크기(pt) 반환."""
    if not shape.has_text_frame:
        return 0
    max_pt = 0
    for para in shape.text_frame.paragraphs:
        for run in para.runs:
            if run.font.size:
                pt = run.font.size.pt
                if pt > max_pt:
                    max_pt = pt
    return max_pt

def get_anchor(shape):
    """vertical anchor: 'top' / 'middle' / 'bottom'."""
    if not shape.has_text_frame:
        return "top"
    a = shape.text_frame.vertical_anchor
    if a == MSO_ANCHOR.MIDDLE:
        return "middle"
    if a == MSO_ANCHOR.BOTTOM:
        return "bottom"
    return "top"

def text_visual_bounds(shape, x, y, w, h):
    """글자가 실제 차지하는 시각 영역 (x, y, w, h_visual)."""
    pt = get_max_font_size(shape)
    if pt == 0:
        return (x, y, w, h)
    line_count = shape.text_frame.text.count("\n") + 1 if shape.text_frame.text else 1
    text_h = font_pt_to_inch(pt) * line_count

    anchor = get_anchor(shape)
    if anchor == "middle":
        actual_y = y + max(0, (h - text_h) / 2)
    elif anchor == "bottom":
        actual_y = y + max(0, h - text_h)
    else:  # top
        actual_y = y
    actual_h = text_h
    return (x, actual_y, w, actual_h)

def rect_overlap(a, b, slack=0.05):
    ax1, ay1, aw, ah = a
    ax2, ay2 = ax1 + aw, ay1 + ah
    bx1, by1, bw, bh = b
    bx2, by2 = bx1 + bw, by1 + bh
    ox = max(0, min(ax2, bx2) - max(ax1, bx1))
    oy = max(0, min(ay2, by2) - max(ay1, by1))
    if ox <= slack or oy <= slack:
        return 0.0
    overlap_area = ox * oy
    a_area = aw * ah
    return overlap_area / a_area if a_area > 0 else 0.0

def has_solid_fill(shape):
    """도형이 solid fill인지 (배경이 텍스트를 가릴 수 있음)."""
    try:
        if not hasattr(shape, "fill"):
            return False
        from pptx.enum.dml import MSO_FILL
        return shape.fill.type == MSO_FILL.SOLID
    except Exception:
        return False

def analyze(p: Path):
    prs = Presentation(p)
    print(f"=== {p.name} (visual overlap v2) ===")
    print(f"Canvas: {prs.slide_width.inches:.2f}\" x {prs.slide_height.inches:.2f}\"\n")

    total_issues = 0

    for si, slide in enumerate(prs.slides, 1):
        print(f"--- Slide {si} ({len(slide.shapes)} shapes) ---")
        boxes = []
        # 텍스트 없지만 fill 있는 도형(카드 배경 등)도 별도 수집
        bg_shapes = []
        for i, shp in enumerate(slide.shapes):
            try:
                x = emu_to_inch(shp.left)
                y = emu_to_inch(shp.top)
                w = emu_to_inch(shp.width)
                h = emu_to_inch(shp.height)
            except Exception:
                continue
            txt = shape_text(shp)
            box = (x, y, w, h)
            if not txt:
                # 텍스트 없는 도형 — 배경 fill인지 확인
                if has_solid_fill(shp) and w > 0.5 and h > 0.2:
                    bg_shapes.append({"i": i, "box": box})
                continue
            visual = text_visual_bounds(shp, x, y, w, h)
            pt = get_max_font_size(shp)
            anchor = get_anchor(shp)
            boxes.append({"i": i, "txt": txt, "box": box, "visual": visual, "pt": pt, "anchor": anchor})

        # 텍스트가 배경 도형에 가려지는 경우 (later z-order 도형이 글자를 덮음)
        # PPT 도형 추가 순서 = z-order. 나중에 추가된 도형이 위에 옴.
        # → 텍스트 박스 인덱스보다 큰 인덱스의 배경 도형이 그 텍스트를 덮음
        bg_cover = []
        for tb in boxes:
            for bg in bg_shapes:
                if bg["i"] <= tb["i"]:
                    continue  # 배경이 먼저 그려졌으면 텍스트가 위에 옴 (OK)
                ratio = rect_overlap(tb["visual"], bg["box"], slack=0.05)
                if ratio > 0.20:  # 텍스트의 20% 이상이 배경에 덮임
                    bg_cover.append((tb, bg, ratio))

        if bg_cover:
            for tb, bg, ratio in bg_cover:
                print(f"  COVER  text #{tb['i']:2d} '{tb['txt'][:25]}' covered by bg shape #{bg['i']} ({ratio*100:.0f}%)")
                total_issues += 1

        # 박스가 폰트보다 작은 경우 경고
        font_box_warns = []
        for b in boxes:
            box_h = b["box"][3]
            visual_h = b["visual"][3]
            if visual_h > box_h + 0.05 and b["pt"] >= 24:
                font_box_warns.append(b)

        if font_box_warns:
            for w in font_box_warns:
                print(f"  WARN font>box  #{w['i']:2d} '{w['txt'][:30]}' "
                      f"box_h={w['box'][3]:.2f} font_h={w['visual'][3]:.2f} "
                      f"({w['pt']:.0f}pt, {w['anchor']})")
                total_issues += 1

        # 시각 영역 기준 겹침
        v_overlaps = []
        for i, a in enumerate(boxes):
            for j, b in enumerate(boxes):
                if j <= i:
                    continue
                ratio = rect_overlap(a["visual"], b["visual"], slack=0.08)
                # 박스 자체 겹침은 제외하고 시각 영역만 새로 잡힌 케이스 강조
                ratio_box = rect_overlap(a["box"], b["box"], slack=0.08)
                if ratio > 0.05:
                    kind = "VISUAL" if ratio_box <= 0.05 else "BOX"
                    v_overlaps.append((a, b, ratio, kind))

        if v_overlaps:
            print(f"  >> {len(v_overlaps)} OVERLAP(s):")
            for a, b, ratio, kind in v_overlaps:
                print(f"     [{kind}] #{a['i']:2d} '{a['txt'][:22]}' x #{b['i']:2d} '{b['txt'][:22]}' ({ratio*100:.0f}%)")
                total_issues += 1
        else:
            if not font_box_warns:
                print(f"  >> No overlap, no font/box mismatch")
        print()

    print(f"=== Total issues: {total_issues} ===")
    return total_issues

if __name__ == "__main__":
    target = Path(sys.argv[1]) if len(sys.argv) > 1 else \
             Path(__file__).resolve().parent.parent / "exports" / "olist_demo_v1.pptx"
    sys.exit(0 if analyze(target) == 0 else 1)
