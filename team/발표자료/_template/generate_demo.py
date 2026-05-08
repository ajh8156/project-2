"""
Olist 재구매 전략 발표 — 시범 5장 PPTX 생성
디자인 시스템 v1.2 준수

생성 슬라이드:
  Slide 1  — Type 1  Cover
  Slide 4  — Type 10 Executive Summary (SCR)
  Slide 9  — Type 7  Quote/Stat (5점 편향 ★ 발표 정점)
  Slide 22 — Type 6  Timeline (Lifecycle Journey Map ★ 시각 정점)
  Slide 29 — Type 9  Closing (THANK YOU 다크)

출력: team/발표자료/exports/olist_demo_v1.pptx
"""

from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn
from lxml import etree

# ============================================================
# 디자인 토큰 (v1.2 정확값)
# ============================================================
OLIST_BLUE = RGBColor(0x1E, 0x40, 0xFF)
BLUE_TINT  = RGBColor(0xEE, 0xF2, 0xFF)
NEAR_BLACK = RGBColor(0x08, 0x08, 0x08)
GRAY_800   = RGBColor(0x22, 0x22, 0x22)
GRAY_600   = RGBColor(0x5A, 0x5A, 0x5A)
GRAY_400   = RGBColor(0x9C, 0xA3, 0xAF)
GRAY_200   = RGBColor(0xE5, 0xE7, 0xEB)
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)

# Semantic 강조 색상 (v1.3 — 블루 단색 시스템에서 시각 강도 차이용)
# 한 슬라이드 최대 3곳, dot/뱃지/작은 강조 한정 (큰 면적 금지)
CRITICAL = RGBColor(0xEE, 0x1D, 0x36)   # 위기·부정 (Retention 붕괴, 이탈)
WARNING  = RGBColor(0xFF, 0xAE, 0x13)   # 주의·중간 (Activation 흔들림)
SUCCESS  = RGBColor(0x00, 0xD7, 0x22)   # 긍정·달성 (KR 목표 달성)

# 캔버스 (13.333" × 7.5")
SLIDE_W = 13.333
SLIDE_H = 7.5
PAD_X = 0.56
PAD_Y_TOP = 0.35
RIGHT_EDGE = 12.50      # 안전 마진 (우측 정렬 텍스트는 이 안쪽)
FOOTER_Y = 7.15

# Pretendard weight → 폰트 패밀리명 매핑
FONT_BLACK     = "Pretendard Black"
FONT_EXTRABOLD = "Pretendard ExtraBold"
FONT_BOLD      = "Pretendard"           # Bold는 run.bold=True
FONT_SEMIBOLD  = "Pretendard SemiBold"
FONT_MEDIUM    = "Pretendard Medium"
FONT_REGULAR   = "Pretendard"
FONT_LIGHT     = "Pretendard Light"

# ============================================================
# Helper: Text 추가
# ============================================================
def add_text(slide, x, y, w, h, text,
             size=14, font="Pretendard", bold=False,
             color=NEAR_BLACK, align="left",
             char_spacing=0, line_spacing=None,
             anchor="top"):
    """char_spacing: pt 단위 (0.1em ≈ size * 0.1)"""
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    tf.word_wrap = True
    if anchor == "middle":
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    elif anchor == "bottom":
        tf.vertical_anchor = MSO_ANCHOR.BOTTOM

    p = tf.paragraphs[0]
    p.alignment = {"left": PP_ALIGN.LEFT, "center": PP_ALIGN.CENTER, "right": PP_ALIGN.RIGHT}[align]
    if line_spacing:
        p.line_spacing = line_spacing

    run = p.add_run()
    run.text = text
    run.font.name = font
    run.font.size = Pt(size)
    run.font.color.rgb = color
    if bold:
        run.font.bold = True
    if char_spacing > 0:
        # XML 직접 수정 — spc는 1/100 pt 단위
        rPr = run._r.get_or_add_rPr()
        rPr.set("spc", str(int(char_spacing * 100)))
    return tb

# ============================================================
# Helper: 한 박스 안에 여러 색 run (Page Title 색 분리용)
# ============================================================
def add_text_multicolor(slide, x, y, w, h, runs,
                       size=36, font="Pretendard", bold=True,
                       align="left", anchor="top",
                       char_spacing=0, line_spacing=None):
    """runs = [(text, color), ...] — 한 박스 안 여러 색 run."""
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    tf.word_wrap = True
    if anchor == "middle":
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    elif anchor == "bottom":
        tf.vertical_anchor = MSO_ANCHOR.BOTTOM

    p = tf.paragraphs[0]
    p.alignment = {"left": PP_ALIGN.LEFT, "center": PP_ALIGN.CENTER, "right": PP_ALIGN.RIGHT}[align]
    if line_spacing:
        p.line_spacing = line_spacing

    for text, color in runs:
        run = p.add_run()
        run.text = text
        run.font.name = font
        run.font.size = Pt(size)
        run.font.color.rgb = color
        if bold:
            run.font.bold = True
        if char_spacing > 0:
            rPr = run._r.get_or_add_rPr()
            rPr.set("spc", str(int(char_spacing * 100)))
    return tb

# ============================================================
# Helper: 직사각형 (배경/카드/색띠)
# ============================================================
def add_rect(slide, x, y, w, h, fill=WHITE, line=None, line_width=0.75):
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line
        shp.line.width = Pt(line_width)
    shp.shadow.inherit = False
    return shp

# ============================================================
# Helper: 가로 구분선 (h: 0.01 — PowerPoint 호환)
# ============================================================
def add_h_line(slide, x, y, w, color=GRAY_200):
    return add_rect(slide, x, y, w, 0.01, fill=color)

# ============================================================
# Helper: 원형 (Journey Map 노드, Persona Avatar)
# ============================================================
def add_circle(slide, cx, cy, d, fill=OLIST_BLUE, line=None):
    """cx, cy = 중심 / d = 지름"""
    x = cx - d/2
    y = cy - d/2
    shp = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x), Inches(y), Inches(d), Inches(d))
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line
        shp.line.width = Pt(1)
    shp.shadow.inherit = False
    return shp

# ============================================================
# 공통 Header (모든 슬라이드)
# ============================================================
def add_header(slide, active_idx=-1, dark=False):
    """active_idx: 0=Diagnose, 1=Bottleneck, 2=Solution, 3=Outcome, -1=none"""
    # olist 로고
    logo_color = WHITE if dark else OLIST_BLUE
    add_text(slide, PAD_X, PAD_Y_TOP, 1.5, 0.35,
             "olist", size=16, font=FONT_BLACK, color=logo_color,
             char_spacing=-0.5)

    # 4 nav (우측 정렬)
    navs = ["DIAGNOSE", "BOTTLENECK", "SOLUTION", "OUTCOME"]
    nav_w = 1.15  # 각 메뉴 폭
    gap = 0.05
    nav_total = len(navs) * nav_w + (len(navs)-1) * gap
    nav_start_x = RIGHT_EDGE - nav_total

    for i, item in enumerate(navs):
        active = (i == active_idx)
        if dark:
            col = WHITE if active else RGBColor(0xCC, 0xD3, 0xFF)
        else:
            col = OLIST_BLUE if active else GRAY_400
        weight_font = FONT_BOLD if active else FONT_MEDIUM
        # charSpacing: 다크 모드에선 0 (v1.2 §11.3)
        cs = 0 if dark else 0.9  # 0.1em @ 9pt ≈ 0.9 pt
        add_text(slide,
                 nav_start_x + i*(nav_w + gap), PAD_Y_TOP + 0.04,
                 nav_w, 0.30,
                 item, size=9, font=weight_font, bold=active,
                 color=col, char_spacing=cs, align="right")

# ============================================================
# 공통 Footer (모든 슬라이드)
# ============================================================
def add_footer(slide, page_num, dark=False):
    if dark:
        col = RGBColor(0xCC, 0xD3, 0xFF)
        cs = 0
    else:
        col = GRAY_400
        cs = 0.7
    add_text(slide, PAD_X, FOOTER_Y, 6.5, 0.25,
             "NEVER—ENDING · OLIST 재구매 전략",
             size=7, font=FONT_REGULAR, color=col, char_spacing=cs)
    add_text(slide, RIGHT_EDGE - 2, FOOTER_Y, 2, 0.25,
             f"2018 · {page_num:02d}",
             size=7, font=FONT_REGULAR, color=col,
             char_spacing=cs, align="right")

# ============================================================
# Slide 1 — Cover (Type 1)
# ============================================================
def make_slide_1(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    add_header(s, active_idx=-1)

    # Mega "OLIST" 중앙 — 박스 폰트 높이보다 크게 + middle anchor
    add_text(s, 0, 2.0, SLIDE_W, 2.2, "OLIST",
             size=140, font=FONT_BLACK, color=OLIST_BLUE,
             align="center", char_spacing=-5.6, anchor="middle")

    # Meta (기간 · 팀명)
    add_text(s, 0, 4.5, SLIDE_W, 0.35,
             "2017.01 — 2018.08  ·  데2터로말해조",
             size=11, font=FONT_BOLD, bold=True, color=NEAR_BLACK,
             align="center")

    # Lead
    add_text(s, 0, 5.0, SLIDE_W, 0.5,
             "올리스트 성장 진단 컨설팅 보고서. 이탈률 개선을 위한 액션플랜.",
             size=11, font=FONT_REGULAR, color=GRAY_600,
             align="center")

    add_footer(s, page_num=1)
    return s

# ============================================================
# Slide 4 — Executive Summary (Type 10, SCR 3카드)
# ============================================================
def make_slide_4(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_header(s, active_idx=0)  # Diagnose

    # Page Title — 한 박스 + run 색 분리 (이전 두 박스 겹침 패치)
    add_text_multicolor(
        s, PAD_X, 1.0, 12.21, 0.7,
        runs=[("NO REPURCHASE", OLIST_BLUE),
              ("  STRUCTURE",   NEAR_BLACK)],
        size=36, font=FONT_BLACK, bold=True)
    # Lead
    add_text(s, PAD_X, 1.65, 11, 0.4,
             "Olist 재구매 진단 결과 — 한 장 요약. 100명 중 0.5명만 돌아옵니다.",
             size=11, font=FONT_REGULAR, color=GRAY_600)

    # 3 Cards (SCR)
    card_y = 2.4
    card_h = 4.0
    gap = 0.2
    card_w = (12.21 - 2*gap) / 3  # 약 3.94"
    cards = [
        {"label": "SITUATION", "label_color": GRAY_400,
         "head": "재구매율 3.00%",
         "body": "100명 중 97명이 한 번 사고 돌아오지 않습니다.\n매출 94%가 일회성 고객 의존.",
         "stat_label": "재구매 주문 비중", "stat_value": "6.13%",
         "stat_color": GRAY_800, "border_color": GRAY_400, "bg": WHITE},
        {"label": "COMPLICATION", "label_color": GRAY_800,
         "head": "BF cohort 0.56%",
         "body": "작년 BF로 만든 고객조차 30일 내 재구매율이 0.56%.\n마케팅 비용 회수 불가.",
         "stat_label": "BF 30일 재구매율", "stat_value": "0.56%",
         "stat_color": GRAY_800, "border_color": GRAY_800, "bg": WHITE},
        {"label": "RESOLUTION", "label_color": OLIST_BLUE,
         "head": "TF 3개로 재구매 구조를 만든다",
         "body": "배송 품질 + 첫 경험 + 재구매 유도 — 3개 TF로\nKR1 재구매율 4.5% 달성.",
         "stat_label": "1년 목표 KR1", "stat_value": "4.5%",
         "stat_color": OLIST_BLUE, "border_color": OLIST_BLUE, "bg": BLUE_TINT},
    ]
    for i, c in enumerate(cards):
        cx = PAD_X + i*(card_w + gap)
        # 배경 카드
        add_rect(s, cx, card_y, card_w, card_h, fill=c["bg"])
        # 좌측 6px 색띠
        add_rect(s, cx, card_y, 0.05, card_h, fill=c["border_color"])
        # 라벨 (Meta UPPER)
        add_text(s, cx + 0.3, card_y + 0.3, card_w - 0.6, 0.3,
                 c["label"], size=8, font=FONT_MEDIUM,
                 color=c["label_color"], char_spacing=1.2)
        # Head
        add_text(s, cx + 0.3, card_y + 0.9, card_w - 0.6, 0.7,
                 c["head"], size=18, font=FONT_BOLD, bold=True,
                 color=NEAR_BLACK, line_spacing=1.2)
        # Body
        add_text(s, cx + 0.3, card_y + 1.9, card_w - 0.6, 1.2,
                 c["body"], size=10, font=FONT_REGULAR,
                 color=GRAY_600, line_spacing=1.4)
        # Divider
        add_h_line(s, cx + 0.3, card_y + 3.1, card_w - 0.6)
        # Stat label
        add_text(s, cx + 0.3, card_y + 3.25, card_w - 0.6, 0.25,
                 c["stat_label"], size=8, font=FONT_MEDIUM,
                 color=GRAY_600, char_spacing=1.0)
        # Stat value
        add_text(s, cx + 0.3, card_y + 3.5, card_w - 0.6, 0.5,
                 c["stat_value"], size=32, font=FONT_BLACK, bold=True,
                 color=c["stat_color"], char_spacing=-0.6)

    # Bottom: Read More
    add_h_line(s, PAD_X, 6.65, 12.21)
    add_text(s, PAD_X, 6.78, 1.5, 0.25, "READ MORE",
             size=8, font=FONT_MEDIUM, color=GRAY_600, char_spacing=1.2)
    add_text(s, PAD_X + 1.6, 6.78, 10, 0.25,
             "Diagnose → Section 01  ·  Bottleneck → Section 02  ·  Solution → Section 03",
             size=9, font=FONT_REGULAR, color=GRAY_600)

    add_footer(s, page_num=4)
    return s

# ============================================================
# Slide 9 — Quote/Stat (Type 7) ★ 발표 정점
# ============================================================
def make_slide_9(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_header(s, active_idx=1)  # Bottleneck

    # Top Meta — 위로 이동
    add_text(s, 0, 1.2, SLIDE_W, 0.3,
             "CORE INSIGHT  ·  RATING × REPURCHASE",
             size=9, font=FONT_MEDIUM, color=GRAY_600,
             char_spacing=1.4, align="center")

    # Mega Number "3.12%" — 폰트 180pt(박스 3.0"에 안전 마진 확보) + middle anchor
    add_text(s, 0, 1.5, SLIDE_W, 3.0, "3.12%",
             size=180, font=FONT_BLACK, color=OLIST_BLUE,
             align="center", char_spacing=-7.2, anchor="middle")

    # Headline — 위치 조정
    add_text(s, 0, 4.65, SLIDE_W, 0.55,
             "5점 리뷰만 재구매를 만든다",
             size=24, font=FONT_BLACK, bold=True, color=NEAR_BLACK,
             align="center")

    # Body
    add_text(s, 1.5, 5.3, SLIDE_W - 3, 0.9,
             "1점·2점·3점·4점 리뷰 모두 재구매율 ~2.5% 동일. 5점만 3.12%.\n5점을 만드는 것은 정시 배송(3.04%) — 지연 시 2.51%.",
             size=12, font=FONT_REGULAR, color=GRAY_600,
             align="center", line_spacing=1.5)

    # Bottom Quote
    add_text(s, 0, 6.35, SLIDE_W, 0.4,
             "재구매 구조의 뿌리는 배송 경험이다.",
             size=14, font=FONT_BOLD, bold=True, color=NEAR_BLACK,
             align="center")

    # Source caption
    add_text(s, 0, 6.85, SLIDE_W, 0.25,
             "SOURCE: Olist Reviews × Repurchase Linkage  ·  N=99,441",
             size=8, font=FONT_REGULAR, color=GRAY_400,
             char_spacing=1.0, align="center")

    add_footer(s, page_num=9)
    return s

# ============================================================
# Slide 22 — Lifecycle Journey Map (Type 6) ★ 시각 정점
# (재설계: Y 좌표 분리, Page Title 단일 박스, 라벨 폭 1.6)
# ============================================================
def make_slide_22(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_header(s, active_idx=2)  # Solution

    # Page Title — 단일 박스 + run 색 분리 (폭 7.5로 축소: 우측 Meta와 겹침 패치)
    add_text_multicolor(
        s, PAD_X, 1.0, 7.5, 0.7,
        runs=[("LIFECYCLE", OLIST_BLUE),
              ("  JOURNEY MAP", NEAR_BLACK)],
        size=36, font=FONT_BLACK, bold=True)

    # Meta (우측)
    add_text(s, RIGHT_EDGE - 4, 1.20, 4, 0.3,
             "7 TOUCHPOINTS  ·  D+0 → D+90",
             size=9, font=FONT_MEDIUM, color=GRAY_600,
             char_spacing=1.4, align="right")

    # ── 좌표 재설계 (Y 분리) ──────────────────
    LABEL_TOP_Y  = 1.95   # 노드 상단 라벨
    NODE_CY      = 2.95   # 원형 노드 중심 (가로 라인도 여기)
    LABEL_BOT_Y  = 3.45   # 노드 하단 라벨
    STAGE_Y      = 4.20   # Stage Labels (그룹 띠)
    TEAM_Y       = 4.65   # Team Labels
    HIGHLIGHT_Y  = 5.55   # Bottom Highlight Box

    LABEL_W = 1.6  # 노드 라벨 박스 폭 (이전 2.0 → 좌우 겹침 패치)

    # 노드 정의 — Semantic 색대 그라데이션:
    # 첫경험(SUCCESS) → 재구매(OLIST_BLUE) → 이탈방지(WARNING→CRITICAL)
    nodes = [
        {"label": "첫구매", "sub": "BUY",         "color": SUCCESS,    "filled": True},
        {"label": "D+0",    "sub": "Thank-you",   "color": SUCCESS,    "filled": False, "tint": True},
        {"label": "D+3",    "sub": "만족도 체크",   "color": OLIST_BLUE, "filled": False, "tint": True},
        {"label": "D+7",    "sub": "앱푸시",       "color": OLIST_BLUE, "filled": True},
        {"label": "D+30",   "sub": "카카오·이메일", "color": OLIST_BLUE, "filled": True},
        {"label": "D+60",   "sub": "카카오 시크릿", "color": WARNING,    "filled": True},
        {"label": "D+90",   "sub": "윈백 캠페인",   "color": CRITICAL,   "filled": True},
    ]
    n = len(nodes)
    node_d = 0.6
    span_left = 1.2
    span_right = 12.2
    step = (span_right - span_left) / (n - 1)

    # 가로 라인 (노드 중심 통과)
    add_h_line(s, 1.0, NODE_CY, 11.5, color=GRAY_200)

    # 7개 노드
    for i, nd in enumerate(nodes):
        cx = span_left + i * step
        if nd["filled"]:
            add_circle(s, cx, NODE_CY, node_d, fill=nd["color"])
            num_color = WHITE
        elif nd.get("tint"):
            add_circle(s, cx, NODE_CY, node_d, fill=BLUE_TINT, line=OLIST_BLUE)
            num_color = OLIST_BLUE
        else:
            add_circle(s, cx, NODE_CY, node_d, fill=WHITE, line=GRAY_400)
            num_color = GRAY_600
        # Number 0X
        add_text(s, cx - node_d/2, NODE_CY - 0.13, node_d, 0.3,
                 f"{i+1:02d}", size=11, font=FONT_BOLD, bold=True,
                 color=num_color, align="center")
        # 노드 상단 라벨
        add_text(s, cx - LABEL_W/2, LABEL_TOP_Y, LABEL_W, 0.3, nd["label"],
                 size=12, font=FONT_BOLD, bold=True, color=NEAR_BLACK, align="center")
        # 노드 하단 라벨
        add_text(s, cx - LABEL_W/2, LABEL_BOT_Y, LABEL_W, 0.3, nd["sub"],
                 size=9, font=FONT_REGULAR, color=GRAY_600, align="center")

    # ── Stage / Team 라벨: 노드 그룹 중심 정렬, 폭 2.5 고정 ──
    GROUP_W = 2.5
    groups = [
        # First Experience (SUCCESS): 노드 0,1,2
        {"stage": "FIRST EXPERIENCE", "team": "Product · 물류 → CRM",
         "from": 0, "to": 2, "stage_color": SUCCESS, "team_color": GRAY_600},
        # Repurchase Trigger (OLIST_BLUE): 노드 3,4
        {"stage": "REPURCHASE TRIGGER", "team": "CRM · MD",
         "from": 3, "to": 4, "stage_color": OLIST_BLUE, "team_color": OLIST_BLUE},
        # Churn Prevention (WARNING→CRITICAL): 노드 5,6
        {"stage": "CHURN PREVENTION", "team": "CRM · Product",
         "from": 5, "to": 6, "stage_color": CRITICAL, "team_color": GRAY_600},
    ]
    for g in groups:
        cx_start = span_left + g["from"] * step
        cx_end   = span_left + g["to"]   * step
        cx_center = (cx_start + cx_end) / 2
        x = cx_center - GROUP_W / 2
        # Stage Label
        add_text(s, x, STAGE_Y, GROUP_W, 0.3, g["stage"],
                 size=9, font=FONT_MEDIUM, color=g["stage_color"],
                 char_spacing=1.4, align="center")
        # Team Label
        add_text(s, x, TEAM_Y, GROUP_W, 0.3, g["team"],
                 size=9, font=FONT_MEDIUM, color=g["team_color"],
                 char_spacing=1.0, align="center")

    # Bottom Highlight Box
    add_rect(s, PAD_X, HIGHLIGHT_Y, 12.21, 1.0, fill=BLUE_TINT)
    add_rect(s, PAD_X, HIGHLIGHT_Y, 0.05, 1.0, fill=OLIST_BLUE)
    add_text(s, PAD_X + 0.3, HIGHLIGHT_Y + 0.18, 1.5, 0.3, "★ KEY",
             size=8, font=FONT_MEDIUM, color=OLIST_BLUE, char_spacing=1.4)
    add_text(s, PAD_X + 0.3, HIGHLIGHT_Y + 0.5, 11.7, 0.5,
             "이게 우리가 만들 7터치포인트 라이프사이클 — D+0부터 D+90까지 자동으로 굴러갑니다.",
             size=14, font=FONT_BOLD, bold=True, color=NEAR_BLACK)

    add_footer(s, page_num=22)
    return s

# ============================================================
# Slide 29 — THANK YOU (Type 9, 다크 풀블리드)
# ============================================================
def make_slide_29(prs, page_num=29):
    s = prs.slides.add_slide(prs.slide_layouts[6])

    # 풀블리드 블루 배경
    bg = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0,
                              Inches(SLIDE_W), Inches(SLIDE_H))
    bg.fill.solid()
    bg.fill.fore_color.rgb = OLIST_BLUE
    bg.line.fill.background()
    bg.shadow.inherit = False

    add_header(s, active_idx=-1, dark=True)

    # Mega "THANK YOU" — 박스 크게 + middle anchor + 다크 charSpacing 0 (v1.2 §11.3)
    add_text(s, 0, 2.2, SLIDE_W, 2.5, "THANK YOU",
             size=160, font=FONT_BLACK, color=WHITE,
             align="center", char_spacing=0, anchor="middle")

    # Lead — 위치 조정
    add_text(s, 0, 4.95, SLIDE_W, 0.4,
             "Never-ending stories start here.",
             size=14, font=FONT_LIGHT, color=RGBColor(0xCC, 0xD3, 0xFF),
             align="center")

    # Bottom Meta
    add_text(s, 0, 5.85, SLIDE_W, 0.3,
             "데2터로말해조  ·  6인 외부 컨설팅",
             size=10, font=FONT_MEDIUM, color=RGBColor(0xCC, 0xD3, 0xFF),
             align="center")

    add_footer(s, page_num=page_num, dark=True)
    return s

# ============================================================
# Main
# ============================================================
def main():
    out_dir = Path(__file__).resolve().parent.parent / "exports"
    out_dir.mkdir(exist_ok=True)
    out_path = out_dir / "olist_demo_v1.pptx"

    prs = Presentation()
    prs.slide_width = Inches(SLIDE_W)
    prs.slide_height = Inches(SLIDE_H)

    print("Generating Slide 1 - Cover ...")
    make_slide_1(prs)
    print("Generating Slide 4 - Executive Summary (SCR) ...")
    make_slide_4(prs)
    print("Generating Slide 9 - Quote/Stat (Peak) ...")
    make_slide_9(prs)
    print("Generating Slide 22 - Lifecycle Journey Map (Visual Peak) ...")
    make_slide_22(prs)
    print("Generating Slide 29 - THANK YOU (Dark) ...")
    make_slide_29(prs)

    prs.save(out_path)
    print(f"\nDone -> {out_path}")
    print(f"  size: {out_path.stat().st_size:,} bytes")

if __name__ == "__main__":
    main()
