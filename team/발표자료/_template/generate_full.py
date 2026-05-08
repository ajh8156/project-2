"""
Olist 발표 PPTX 전체 생성기 — Batch A 진행 중 (11장: 시범 5 + 신규 6)
generate_demo.py 의 헬퍼·시범 함수를 import해 재사용.
"""

import sys
from pathlib import Path

sys.path.insert(0, str(Path(__file__).resolve().parent))
from generate_demo import (
    OLIST_BLUE, BLUE_TINT, NEAR_BLACK, GRAY_800, GRAY_600, GRAY_400, GRAY_200, WHITE,
    CRITICAL, WARNING, SUCCESS,
    SLIDE_W, SLIDE_H, PAD_X, PAD_Y_TOP, RIGHT_EDGE, FOOTER_Y,
    FONT_BLACK, FONT_EXTRABOLD, FONT_BOLD, FONT_SEMIBOLD, FONT_MEDIUM, FONT_REGULAR, FONT_LIGHT,
    add_text, add_text_multicolor, add_rect, add_h_line, add_circle,
    add_header, add_footer,
    make_slide_1, make_slide_4, make_slide_9, make_slide_22, make_slide_29,
)
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

# 산출물 베이스 경로 (team/ 디렉토리)
TEAM_DIR = Path(__file__).resolve().parent.parent.parent

def add_picture_safe(slide, rel_path, x, y, w, h):
    """이미지 자동 삽입. 없으면 placeholder."""
    p = TEAM_DIR / rel_path
    if not p.exists():
        add_rect(slide, x, y, w, h, fill=GRAY_200)
        add_text(slide, x, y + h/2 - 0.15, w, 0.3,
                 f"[{p.name}]", size=10, font=FONT_REGULAR,
                 color=GRAY_600, align="center")
        return None
    return slide.shapes.add_picture(str(p), Inches(x), Inches(y), Inches(w), Inches(h))


# ============================================================
# Slide 2 — Type 4 변형: 의뢰 이메일 카드
# ============================================================
def make_slide_2(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_header(s, active_idx=0)  # Diagnose

    # 좌측 멀티라인 mega title
    add_text(s, PAD_X, 1.5, 6.5, 1.0, "AN",
             size=60, font=FONT_BLACK, bold=True, color=NEAR_BLACK)
    add_text(s, PAD_X, 2.4, 6.5, 1.0, "EMAIL",
             size=60, font=FONT_BLACK, bold=True, color=OLIST_BLUE)
    add_text(s, PAD_X, 3.3, 6.5, 1.0, "FROM OLIST",
             size=60, font=FONT_BLACK, bold=True, color=NEAR_BLACK)

    # 좌측 하단 Lead
    add_text(s, PAD_X, 5.0, 6.0, 0.9,
             "2018.04.10 — Olist Growth Director Carlos Silva로부터\n의뢰가 도착했습니다.",
             size=12, font=FONT_REGULAR, color=GRAY_600, line_spacing=1.5)

    # 우측 이메일 카드
    EX, EY, EW, EH = 7.0, 1.0, 5.5, 5.6
    add_rect(s, EX, EY, EW, EH, fill=WHITE, line=GRAY_200, line_width=0.5)

    # 메일함 헤더 띠
    add_rect(s, EX, EY, EW, 0.4, fill=GRAY_200)
    add_text(s, EX + 0.2, EY + 0.08, EW - 0.4, 0.3,
             "받은 메일함  /  2018.04.10 (火) 14:23",
             size=8, font=FONT_MEDIUM, color=GRAY_600, char_spacing=1.0)

    # From/To/제목
    iy = EY + 0.55
    add_text(s, EX + 0.2, iy, 0.7, 0.3, "From:",
             size=9, font=FONT_REGULAR, color=GRAY_600)
    add_text(s, EX + 0.95, iy, EW - 1.15, 0.3,
             "Carlos Silva <carlos.silva@olist.com>",
             size=10, font=FONT_BOLD, bold=True, color=NEAR_BLACK)
    add_text(s, EX + 0.95, iy + 0.25, EW - 1.15, 0.3,
             "Olist Growth Director",
             size=8, font=FONT_REGULAR, color=GRAY_600)

    iy += 0.6
    add_text(s, EX + 0.2, iy, 0.7, 0.3, "To:",
             size=9, font=FONT_REGULAR, color=GRAY_600)
    add_text(s, EX + 0.95, iy, EW - 1.15, 0.3,
             "데2터로말해조 컨설팅",
             size=10, font=FONT_REGULAR, color=NEAR_BLACK)

    iy += 0.32
    add_text(s, EX + 0.2, iy, 0.7, 0.3, "제목:",
             size=9, font=FONT_REGULAR, color=GRAY_600)
    add_text(s, EX + 0.95, iy, EW - 1.15, 0.3,
             "[의뢰] 성장 정체 진단 — 이탈률 개선의 건",
             size=10, font=FONT_BOLD, bold=True, color=NEAR_BLACK)

    # 본문 영역 구분선
    iy += 0.4
    add_h_line(s, EX + 0.2, iy, EW - 0.4)

    # 본문
    iy += 0.15
    body = ("안녕하세요. 올리스트 Growth Team의\n"
            "Carlos Silva입니다.\n"
            "\n"
            "최근 1년간 신규 가입자는 꾸준히 늘었지만,\n"
            "매출 성장이 눈에 띄게 둔화되었습니다.\n"
            "특히 한 번 구매하고 돌아오지 않는 고객이\n"
            "너무 많습니다.")
    add_text(s, EX + 0.2, iy, EW - 0.4, 1.7, body,
             size=10, font=FONT_REGULAR, color=NEAR_BLACK, line_spacing=1.4)

    # 인용 강조
    iy += 1.85
    add_text(s, EX + 0.2, iy, EW - 0.4, 0.4,
             '"성장이 정체되고 있는 우리를 진단해주세요."',
             size=13, font=FONT_BOLD, bold=True, color=OLIST_BLUE)

    # 서명
    iy += 0.55
    add_text(s, EX + 0.2, iy, EW - 0.4, 0.25, "- Carlos Silva",
             size=9, font=FONT_REGULAR, color=NEAR_BLACK)
    add_text(s, EX + 0.2, iy + 0.22, EW - 0.4, 0.25,
             "Growth Director, Olist",
             size=8, font=FONT_REGULAR, color=GRAY_600)

    add_footer(s, page_num=2)
    return s


# ============================================================
# Slide 3 — Type 11 변형: 컨설팅팀 4분면
# ============================================================
def make_slide_3(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_header(s, active_idx=0)

    # Page Title
    add_text_multicolor(
        s, PAD_X, 1.0, 12.21, 0.8,
        runs=[("WHO WE ARE — ", NEAR_BLACK),
              ("DATA2-RO MAL-HAE-JO", OLIST_BLUE)],
        size=40, font=FONT_BLACK, bold=True)

    # Lead
    add_text(s, PAD_X, 1.9, 12.21, 0.4,
             "외부 데이터 컨설팅팀 6인이 의뢰를 측정 가능한 KPI로 정제했습니다.",
             size=12, font=FONT_REGULAR, color=GRAY_600)

    # 4 Cards
    card_y = 2.7
    card_h = 3.6
    gap = 0.2
    card_w = (12.21 - 3 * gap) / 4
    cards = [
        {"num": "01", "title": "TEAM",
         "desc": "기획자 중심 6인", "tools_label": "ROLE",
         "tools": "외부 컨설팅 — 독립 분석", "highlight": False},
        {"num": "02", "title": "TOOLS",
         "desc": "BigQuery · Looker\nPython", "tools_label": "STACK",
         "tools": "Pandas · matplotlib", "highlight": False},
        {"num": "03", "title": "DATASET",
         "desc": "Olist BR E-Commerce\n+ Marketing Funnel", "tools_label": "SCOPE",
         "tools": "11 tables · 99,441 orders", "highlight": False},
        {"num": "04", "title": "MISSION",
         "desc": "받은 일\n이탈률 개선\n──────\n제안 목표\nKR1 재구매율\n3.0% → 4.5%",
         "tools_label": "", "tools": "", "highlight": True},
    ]
    for i, c in enumerate(cards):
        cx = PAD_X + i * (card_w + gap)
        bg = BLUE_TINT if c["highlight"] else GRAY_200
        add_rect(s, cx, card_y, card_w, card_h, fill=bg)
        if c["highlight"]:
            add_rect(s, cx, card_y, 0.05, card_h, fill=OLIST_BLUE)
        # Number
        num_color = OLIST_BLUE if c["highlight"] else GRAY_600
        add_text(s, cx + 0.25, card_y + 0.3, card_w - 0.5, 0.6,
                 c["num"], size=32, font=FONT_BLACK, bold=True, color=num_color)
        # Title
        add_text(s, cx + 0.25, card_y + 1.05, card_w - 0.5, 0.4,
                 c["title"], size=18, font=FONT_BOLD, bold=True, color=NEAR_BLACK)
        # Desc
        text_color = NEAR_BLACK if c["highlight"] else GRAY_600
        add_text(s, cx + 0.25, card_y + 1.55, card_w - 0.5, 1.8,
                 c["desc"], size=11, font=FONT_REGULAR, color=text_color,
                 line_spacing=1.5)
        if c["tools_label"]:
            add_text(s, cx + 0.25, card_y + card_h - 0.7, card_w - 0.5, 0.25,
                     c["tools_label"], size=8, font=FONT_MEDIUM,
                     color=GRAY_600, char_spacing=1.2)
            add_text(s, cx + 0.25, card_y + card_h - 0.45, card_w - 0.5, 0.3,
                     c["tools"], size=10, font=FONT_REGULAR, color=NEAR_BLACK)

    # Bottom Meta
    add_h_line(s, PAD_X, 6.55, 12.21)
    add_text(s, PAD_X, 6.7, 1.5, 0.25, "DATA SCOPE",
             size=8, font=FONT_MEDIUM, color=GRAY_600, char_spacing=1.2)
    add_text(s, PAD_X + 1.6, 6.7, 10, 0.25,
             "2017.01 — 2018.08  ·  100,000 orders",
             size=10, font=FONT_REGULAR, color=GRAY_600)

    add_footer(s, page_num=3)
    return s


# ============================================================
# Slide 5 — Type 4 변형: 3축 신호등
# ============================================================
def make_slide_5(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_header(s, active_idx=0)

    # Page Title
    add_text_multicolor(
        s, PAD_X, 1.0, 12.21, 0.8,
        runs=[("THREE  ", NEAR_BLACK),
              ("SIGNALS", OLIST_BLUE)],
        size=48, font=FONT_BLACK, bold=True)

    # Lead
    add_text(s, PAD_X, 2.0, 12.21, 0.4,
             "AARRR 3축으로 본 Olist의 현재. 숫자 3개로 끝납니다.",
             size=12, font=FONT_REGULAR, color=GRAY_600)

    # 3 Cards
    card_y = 2.75
    card_h = 3.5
    gap = 0.25
    card_w = (12.21 - 2 * gap) / 3
    cards = [
        # Semantic 색 적용: Acquisition=SUCCESS / Activation=WARNING / Retention=CRITICAL
        {"label": "ACQUISITION", "dot": SUCCESS, "bg": GRAY_200,
         "stat": "93,104", "stat_unit": "명",
         "head": "규모는 충분하다",
         "body": "신규 유입 정상 — 이게 문제가 아닙니다.",
         "highlight": False, "stat_size": 44},
        {"label": "ACTIVATION", "dot": WARNING, "bg": GRAY_200,
         "stat": "8.16% / 12.82%", "stat_unit": "지연 / 저평점",
         "head": "첫 경험이 흔들린다",
         "body": "배송 지연·저평점이 첫 인상을 망칩니다.",
         "highlight": False, "stat_size": 28},
        {"label": "RETENTION", "dot": CRITICAL, "bg": BLUE_TINT,
         "stat": "3.0%", "stat_unit": "재구매율",
         "head": "돌아오지 않는다",
         "body": "진짜 병목 — 100명 중 97명이 이탈합니다.",
         "highlight": True, "stat_size": 44},
    ]
    for i, c in enumerate(cards):
        cx = PAD_X + i * (card_w + gap)
        add_rect(s, cx, card_y, card_w, card_h, fill=c["bg"])
        if c["highlight"]:
            add_rect(s, cx, card_y, 0.05, card_h, fill=OLIST_BLUE)
        # Dot + Label
        add_circle(s, cx + 0.45, card_y + 0.55, 0.18, fill=c["dot"])
        add_text(s, cx + 0.7, card_y + 0.45, card_w - 0.9, 0.3,
                 c["label"], size=10, font=FONT_MEDIUM,
                 color=c["dot"], char_spacing=1.4)
        # Stat
        stat_color = OLIST_BLUE if c["highlight"] else NEAR_BLACK
        add_text(s, cx + 0.3, card_y + 1.05, card_w - 0.6, 0.9,
                 c["stat"], size=c["stat_size"],
                 font=FONT_BLACK, bold=True, color=stat_color, char_spacing=-1)
        # Stat unit
        add_text(s, cx + 0.3, card_y + 2.0, card_w - 0.6, 0.25,
                 c["stat_unit"], size=10, font=FONT_REGULAR, color=GRAY_600)
        # Head
        add_text(s, cx + 0.3, card_y + 2.4, card_w - 0.6, 0.4,
                 c["head"], size=14, font=FONT_BOLD, bold=True, color=NEAR_BLACK)
        # Body
        add_text(s, cx + 0.3, card_y + 2.85, card_w - 0.6, 0.5,
                 c["body"], size=10, font=FONT_REGULAR, color=GRAY_600,
                 line_spacing=1.4)

    # Bottom Highlight Box
    add_rect(s, PAD_X, 6.45, 12.21, 0.7, fill=BLUE_TINT)
    add_rect(s, PAD_X, 6.45, 0.05, 0.7, fill=OLIST_BLUE)
    add_text(s, PAD_X + 0.3, 6.6, 12, 0.4,
             "들어오는 건 잘 되고, 첫 경험이 흔들리고, 다시 안 옵니다.",
             size=14, font=FONT_BOLD, bold=True, color=NEAR_BLACK)

    add_footer(s, page_num=5)
    return s


# ============================================================
# Slide 6 — Type 8 Chart: RFM 4분면 (placeholder + 인사이트)
# ============================================================
def make_slide_6(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_header(s, active_idx=0)

    # Page Title
    add_text_multicolor(
        s, PAD_X, 1.0, 8.5, 0.8,
        runs=[("RFM  ", NEAR_BLACK),
              ("SEGMENTATION", OLIST_BLUE)],
        size=44, font=FONT_BLACK, bold=True)
    # Meta
    add_text(s, RIGHT_EDGE - 4, 1.20, 4, 0.3,
             "RECENCY × FREQUENCY  ·  N=93,104",
             size=9, font=FONT_MEDIUM, color=GRAY_600,
             char_spacing=1.4, align="right")

    # Lead
    add_text(s, PAD_X, 2.0, 12.21, 0.4,
             "이탈재구매 2,562명 — 우리가 회복할 수 있는 우선 타겟.",
             size=12, font=FONT_REGULAR, color=GRAY_600)

    # 좌측 4분면 차트 — Python으로 export한 실제 산점도 PNG 삽입
    CX, CY, CW, CH = PAD_X, 2.75, 7.5, 3.8
    rfm_path = Path(__file__).resolve().parent.parent / "exports" / "rfm-scatter.png"
    if rfm_path.exists():
        # 절대 경로로 직접 add_picture
        s.shapes.add_picture(str(rfm_path), Inches(CX), Inches(CY),
                             Inches(CW), Inches(CH))
    else:
        add_rect(s, CX, CY, CW, CH, fill=GRAY_200)
        add_text(s, CX, CY + CH/2 - 0.15, CW, 0.3,
                 "[RFM 산점도 — generate_rfm_chart.py 실행 필요]",
                 size=12, font=FONT_BOLD, bold=True, color=GRAY_600, align="center")

    # 우측 인사이트 카드
    IX, IY, IW, IH = 8.4, 2.75, 4.13, 3.8
    add_rect(s, IX, IY, IW, IH, fill=BLUE_TINT)
    add_rect(s, IX, IY, 0.05, IH, fill=OLIST_BLUE)
    add_text(s, IX + 0.3, IY + 0.3, IW - 0.6, 0.3,
             "INSIGHT", size=10, font=FONT_MEDIUM,
             color=OLIST_BLUE, char_spacing=1.4)
    add_text(s, IX + 0.3, IY + 0.85, IW - 0.6, 1.2,
             "이탈재구매\n2,562명",
             size=28, font=FONT_BLACK, bold=True, color=OLIST_BLUE,
             line_spacing=1.1)
    add_h_line(s, IX + 0.3, IY + 2.25, IW - 0.6)
    add_text(s, IX + 0.3, IY + 2.45, IW - 0.6, 1.2,
             "한 번 이상 구매했지만\n최근 활동이 없는 고객.\n\n10% 전환만 해도\n+256명 충성 고객 확보.",
             size=11, font=FONT_REGULAR, color=NEAR_BLACK,
             line_spacing=1.5)

    # 하단 출처
    add_text(s, PAD_X, 6.7, 12.21, 0.25,
             "DATA SOURCE: Olist Customer × Order Linkage  ·  93,104 customers (2017.01 — 2018.08)",
             size=8, font=FONT_REGULAR, color=GRAY_400, char_spacing=1.0)

    add_footer(s, page_num=6)
    return s


# ============================================================
# Slide 7 — Type 12 Hypothesis Tree (MECE)
# ============================================================
def make_slide_7(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_header(s, active_idx=1)  # Bottleneck

    # Page Title
    add_text_multicolor(
        s, PAD_X, 1.0, 11, 0.8,
        runs=[("BOTTLENECK  ", NEAR_BLACK),
              ("DECOMPOSITION", OLIST_BLUE)],
        size=40, font=FONT_BLACK, bold=True)

    # Lead
    add_text(s, PAD_X, 1.9, 12.21, 0.4,
             "수식 병목과 여정 병목이 같은 지점에서 만난다 — Retention.",
             size=12, font=FONT_REGULAR, color=GRAY_600)

    # Root (좌)
    RX, RY, RW, RH = PAD_X, 3.6, 2.5, 1.2
    add_rect(s, RX, RY, RW, RH, fill=OLIST_BLUE)
    add_text(s, RX + 0.2, RY + 0.2, RW - 0.4, 0.3,
             "WHY?", size=9, font=FONT_MEDIUM,
             color=RGBColor(0xCC, 0xD3, 0xFF), char_spacing=1.4)
    add_text(s, RX + 0.2, RY + 0.5, RW - 0.4, 0.7,
             "GMV 정체의\n핵심 병목",
             size=18, font=FONT_BOLD, bold=True, color=WHITE,
             line_spacing=1.2)

    # Level 1 (3 nodes)
    L1X = 4.0
    L1W = 3.2
    L1_nodes = [
        {"title": "Acquisition", "sub": "정상 ✓", "y": 2.3, "highlight": False},
        {"title": "Activation", "sub": "흔들림 ⚠", "y": 3.8, "highlight": False},
        {"title": "Retention", "sub": "붕괴 🚨", "y": 5.3, "highlight": True},
    ]
    for n in L1_nodes:
        bg = BLUE_TINT if n["highlight"] else WHITE
        line = OLIST_BLUE if n["highlight"] else GRAY_200
        add_rect(s, L1X, n["y"], L1W, 1.0, fill=bg, line=line, line_width=1)
        add_text(s, L1X + 0.2, n["y"] + 0.2, L1W - 0.4, 0.4,
                 n["title"], size=16, font=FONT_BOLD, bold=True,
                 color=OLIST_BLUE if n["highlight"] else NEAR_BLACK)
        add_text(s, L1X + 0.2, n["y"] + 0.6, L1W - 0.4, 0.3,
                 n["sub"], size=11, font=FONT_REGULAR,
                 color=OLIST_BLUE if n["highlight"] else GRAY_600)

    # Level 2 (6 nodes)
    L2X = 7.5
    L2W = 5.0
    L2_nodes = [
        {"text": "1.1   신규 유입 93,104명 — 정상 규모", "y": 2.5, "highlight": False},
        {"text": "2.1   첫 구매 배송 지연 8.16%", "y": 3.4, "highlight": False},
        {"text": "2.2   첫 구매 저평점 12.82%", "y": 4.0, "highlight": False},
        {"text": "3.1   재구매율 3.0%  ★", "y": 4.9, "highlight": True},
        {"text": "3.2   BF cohort 30일 0.56%  ★", "y": 5.5, "highlight": True},
        {"text": "3.3   Frequency 1.0334 — 수식 병목", "y": 6.1, "highlight": False},
    ]
    for n in L2_nodes:
        bg = BLUE_TINT if n["highlight"] else WHITE
        line = OLIST_BLUE if n["highlight"] else GRAY_200
        add_rect(s, L2X, n["y"], L2W, 0.5, fill=bg, line=line, line_width=1)
        color = OLIST_BLUE if n["highlight"] else NEAR_BLACK
        add_text(s, L2X + 0.2, n["y"] + 0.13, L2W - 0.4, 0.3,
                 n["text"],
                 size=11,
                 font=FONT_BOLD if n["highlight"] else FONT_REGULAR,
                 bold=n["highlight"],
                 color=color)

    # MECE Check (하단)
    add_h_line(s, PAD_X, 6.75, 12.21)
    add_text(s, PAD_X, 6.88, 1.5, 0.25, "MECE CHECK",
             size=8, font=FONT_MEDIUM, color=GRAY_600, char_spacing=1.2)
    add_text(s, PAD_X + 1.7, 6.88, 11, 0.25,
             "GMV = Customers × Frequency 수식 병목과 AARRR 여정 병목이 가리키는 지점은 같다 — Retention.",
             size=10, font=FONT_REGULAR, color=GRAY_600)

    add_footer(s, page_num=7)
    return s


# ============================================================
# Slide 8 — Type 8 Chart: 저평점 4팀 분포
# ============================================================
def make_slide_8(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_header(s, active_idx=1)

    # Page Title
    add_text_multicolor(
        s, PAD_X, 1.0, 9.5, 0.8,
        runs=[("FOUR TEAMS, ", OLIST_BLUE),
              ("ONE PROBLEM", NEAR_BLACK)],
        size=42, font=FONT_BLACK, bold=True)
    # Meta
    add_text(s, RIGHT_EDGE - 4, 1.20, 4, 0.3,
             "저평점 원인 분포  ·  N=12,847",
             size=9, font=FONT_MEDIUM, color=GRAY_600,
             char_spacing=1.4, align="right")

    # Lead
    add_text(s, PAD_X, 2.0, 12.21, 0.4,
             "한 팀이 못 푸는 분포 — TF로 묶어야 하는 정당성.",
             size=12, font=FONT_REGULAR, color=GRAY_600)

    # 100% 누적 막대 (가로) — 4구간
    BAR_X, BAR_Y, BAR_W, BAR_H = PAD_X, 3.0, 12.21, 1.3
    segments = [
        {"team": "물류", "label": "미배송 / 배송 지연", "pct": 42, "color": OLIST_BLUE},
        {"team": "CX", "label": "CS 무응답", "pct": 20, "color": GRAY_800},
        {"team": "Seller Ops", "label": "상품 기대 불일치", "pct": 20, "color": GRAY_600},
        {"team": "Product", "label": "UX 문제", "pct": 18, "color": GRAY_400},
    ]
    sx = BAR_X
    for seg in segments:
        seg_w = BAR_W * (seg["pct"] / 100.0)
        add_rect(s, sx, BAR_Y, seg_w, BAR_H, fill=seg["color"])
        add_text(s, sx + 0.2, BAR_Y + 0.18, seg_w - 0.4, 0.45,
                 f"{seg['pct']}%",
                 size=24, font=FONT_BLACK, bold=True, color=WHITE)
        add_text(s, sx + 0.2, BAR_Y + 0.7, seg_w - 0.4, 0.3,
                 seg["team"],
                 size=12, font=FONT_BOLD, bold=True, color=WHITE)
        add_text(s, sx + 0.2, BAR_Y + 0.97, seg_w - 0.4, 0.25,
                 seg["label"],
                 size=9, font=FONT_REGULAR, color=WHITE)
        sx += seg_w

    # 4팀 매핑 표 (좌표 재배치 — Insight Box와 겹침 패치)
    TY = 4.5
    add_text(s, PAD_X, TY, 12.21, 0.3,
             "TEAM RESPONSIBILITY",
             size=10, font=FONT_MEDIUM, color=GRAY_600, char_spacing=1.4)
    add_h_line(s, PAD_X, TY + 0.35, 12.21)

    teams = [
        {"team": "Logistics", "scope": "배송 운영 — 정시 출고 · 도착", "tf": "TF1"},
        {"team": "CX", "scope": "응답 SLA — 48h 대응률 90%", "tf": "TF2 협업"},
        {"team": "Seller Ops", "scope": "상품 정보 정확성 · 셀러 가이드", "tf": "TF1·TF2"},
        {"team": "Product/UX", "scope": "PDP · 결제 · 검색 UX", "tf": "TF2"},
    ]
    for i, t in enumerate(teams):
        ry = TY + 0.55 + i * 0.35  # 행 간격 0.4 → 0.35
        add_text(s, PAD_X, ry, 2.5, 0.3,
                 t["team"], size=11, font=FONT_BOLD, bold=True, color=NEAR_BLACK)
        add_text(s, PAD_X + 2.7, ry, 7.0, 0.3,
                 t["scope"], size=11, font=FONT_REGULAR, color=GRAY_600)
        add_text(s, PAD_X + 9.8, ry, 2.5, 0.3,
                 t["tf"], size=11, font=FONT_BOLD, bold=True, color=OLIST_BLUE)

    # Insight Box (행 4 끝 6.40 → 박스 시작 6.55, 갭 0.15")
    add_rect(s, PAD_X, 6.55, 12.21, 0.5, fill=BLUE_TINT)
    add_rect(s, PAD_X, 6.55, 0.05, 0.5, fill=OLIST_BLUE)
    add_text(s, PAD_X + 0.3, 6.70, 12, 0.3,
             "원인이 4개 팀에 분산 — 한 팀 단독으로 못 푼다. TF 재편의 근거.",
             size=12, font=FONT_BOLD, bold=True, color=NEAR_BLACK)

    add_footer(s, page_num=8)
    return s


# ============================================================
# Batch B 헬퍼: KPI 게이지·Comparison·Milestone
# ============================================================
def add_kpi_gauge(slide, x, y, w, label, current_val, target_val,
                  progress=0.4, color=OLIST_BLUE, big=False):
    """KPI 게이지 카드 (라벨 + 현재값 + 막대 + 목표값)."""
    val_size = 28 if big else 22
    add_text(slide, x, y, w, 0.3, label,
             size=10, font=FONT_MEDIUM, color=GRAY_600, char_spacing=1.4)
    add_text(slide, x, y + 0.35, w / 2, 0.5, f"현재 {current_val}",
             size=val_size, font=FONT_BOLD, bold=True, color=NEAR_BLACK)
    add_text(slide, x + w / 2, y + 0.35, w / 2, 0.5, f"목표 {target_val}",
             size=val_size, font=FONT_BLACK, bold=True, color=color, align="right")
    track_y = y + 1.0
    add_rect(slide, x, track_y, w, 0.06, fill=GRAY_200)
    add_rect(slide, x, track_y, w * progress, 0.06, fill=color)


def add_milestone_strip(slide, x, y, w, items, accent_color=OLIST_BLUE):
    """가로 마일스톤: items=[(period, title, desc), ...]"""
    n = len(items)
    gap = 0.2
    item_w = (w - (n - 1) * gap) / n
    for i, (period, title, desc) in enumerate(items):
        ix = x + i * (item_w + gap)
        add_text(slide, ix, y, item_w, 0.25, period,
                 size=10, font=FONT_MEDIUM, color=accent_color, char_spacing=1.4)
        add_text(slide, ix, y + 0.32, item_w, 0.4, title,
                 size=14, font=FONT_BOLD, bold=True, color=NEAR_BLACK)
        add_text(slide, ix, y + 0.78, item_w, 0.7, desc,
                 size=10, font=FONT_REGULAR, color=GRAY_600, line_spacing=1.4)


def add_comparison_card(slide, x, y, w, h, label, head, body,
                        stat_label="", stat_value="", highlight=False):
    """Before/After 카드"""
    bg = BLUE_TINT if highlight else GRAY_200
    add_rect(slide, x, y, w, h, fill=bg)
    if highlight:
        add_rect(slide, x, y, 0.05, h, fill=OLIST_BLUE)
    label_color = OLIST_BLUE if highlight else GRAY_600
    add_text(slide, x + 0.3, y + 0.3, w - 0.6, 0.3, label,
             size=10, font=FONT_MEDIUM, color=label_color, char_spacing=1.4)
    add_text(slide, x + 0.3, y + 0.8, w - 0.6, 0.7, head,
             size=22, font=FONT_BOLD, bold=True, color=NEAR_BLACK,
             line_spacing=1.2)
    add_text(slide, x + 0.3, y + 1.85, w - 0.6, h - 2.4, body,
             size=11, font=FONT_REGULAR,
             color=NEAR_BLACK if highlight else GRAY_600, line_spacing=1.5)
    if stat_label:
        add_text(slide, x + 0.3, y + h - 1.05, w - 0.6, 0.25,
                 stat_label, size=8, font=FONT_MEDIUM,
                 color=label_color, char_spacing=1.2)
        stat_color = OLIST_BLUE if highlight else NEAR_BLACK
        add_text(slide, x + 0.3, y + h - 0.75, w - 0.6, 0.6,
                 stat_value, size=36, font=FONT_BLACK, bold=True,
                 color=stat_color, char_spacing=-0.7)


def add_page_title(slide, runs, w=12.21, size=42):
    """Page Title 단일 박스 + run 색 분리"""
    add_text_multicolor(slide, PAD_X, 1.0, w, 0.8, runs=runs,
                        size=size, font=FONT_BLACK, bold=True)


def add_lead(slide, text, y=1.95):
    """Page Title 아래 Lead"""
    add_text(slide, PAD_X, y, 12.21, 0.4, text,
             size=12, font=FONT_REGULAR, color=GRAY_600)


# ============================================================
# Slide 10 — Type 4: Finding 통합 표
# ============================================================
def make_slide_10(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_header(s, active_idx=1)  # Bottleneck

    add_page_title(s, runs=[("THREE FINDINGS, ", NEAR_BLACK),
                            ("THREE TFs", OLIST_BLUE)], size=42)
    add_lead(s, "3가지 원인이 보입니다. 그런데 이 3가지는 모두 한 팀에서 못 풉니다.")

    # 표 헤더 (column 좌표 — 박스 안 겹치게 폭 조정)
    HX, HY = PAD_X, 2.7
    # (label, x_offset_from_PAD_X, width)
    cols = [("FINDING", 0.7, 2.0),
            ("CORE EVIDENCE", 2.9, 5.4),
            ("HYPOTHESIS", 8.5, 1.4),
            ("OWNER TF", 10.1, 2.3)]
    for label, dx, cw in cols:
        add_text(s, PAD_X + dx - 0.7, HY, cw, 0.3, label,
                 size=9, font=FONT_MEDIUM, color=GRAY_600, char_spacing=1.4)
    add_h_line(s, PAD_X, HY + 0.4, 12.21)

    rows = [
        {"num": "①", "name": "배송 지연",
         "evidence": "정시 3.04% vs 지연 2.51% (-0.53%p, 3일 변곡점)",
         "h": "H1 ✓", "tf": "TF1", "highlight": False},
        {"num": "②", "name": "5점 편향",
         "evidence": "1~4점 동일 (~2.5%) / 5점만 3.12% 재구매",
         "h": "H2 ✓ (부록)", "tf": "TF2", "highlight": False},
        {"num": "③", "name": "재구매 유도 부재",
         "evidence": "CRM · 추천 · 쿠폰 인프라 전무",
         "h": "H3 ✓", "tf": "TF3 ★", "highlight": True},
    ]
    ry = HY + 0.55
    for r in rows:
        if r["highlight"]:
            add_rect(s, PAD_X, ry - 0.1, 12.21, 0.95, fill=BLUE_TINT)
            add_rect(s, PAD_X, ry - 0.1, 0.05, 0.95, fill=OLIST_BLUE)
        # Finding num + name (폭 줄임: 1.5 → 1.3, 끝점 2.5)
        add_text(s, PAD_X + 0.2, ry, 0.3, 0.4, r["num"],
                 size=20, font=FONT_BLACK, bold=True,
                 color=OLIST_BLUE if r["highlight"] else GRAY_600)
        add_text(s, PAD_X + 0.65, ry + 0.05, 1.7, 0.3, r["name"],
                 size=13, font=FONT_BOLD, bold=True, color=NEAR_BLACK)
        # Evidence (시작 x: 2.1 → 2.9 으로 0.8" 이동, 폭 5.4)
        add_text(s, PAD_X + 2.9, ry + 0.05, 5.4, 0.6, r["evidence"],
                 size=11, font=FONT_REGULAR, color=GRAY_600, line_spacing=1.4)
        # Hypothesis
        add_text(s, PAD_X + 8.5, ry + 0.05, 1.4, 0.3, r["h"],
                 size=11, font=FONT_BOLD, bold=True,
                 color=OLIST_BLUE if r["highlight"] else NEAR_BLACK)
        # TF
        add_text(s, PAD_X + 10.1, ry + 0.05, 2.3, 0.3, r["tf"],
                 size=14, font=FONT_BLACK, bold=True,
                 color=OLIST_BLUE if r["highlight"] else GRAY_600)
        ry += 1.0

    # 하단 Highlight Box
    add_rect(s, PAD_X, 6.5, 12.21, 0.55, fill=BLUE_TINT)
    add_rect(s, PAD_X, 6.5, 0.05, 0.55, fill=OLIST_BLUE)
    add_text(s, PAD_X + 0.3, 6.65, 12, 0.3,
             "각 원인은 각각 다른 TF의 책임 영역. 다음 장부터 TF 정당성과 실행안을 보여드립니다.",
             size=12, font=FONT_BOLD, bold=True, color=NEAR_BLACK)

    add_footer(s, page_num=10)
    return s


# ============================================================
# Slide 11 — Type 4: TF 정당성
# ============================================================
def make_slide_11(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_header(s, active_idx=2)  # Solution

    add_page_title(s, runs=[("WHY  ", NEAR_BLACK),
                            ("TF", OLIST_BLUE),
                            (",  NOT TEAMS", NEAR_BLACK)], size=44)
    add_lead(s, "저평점 원인 4팀 분산 + 재구매 Journey 7터치포인트 분산 = TF 재편 정당성.")

    # 좌측 미니 분포 (슬8 4팀 분산 압축)
    LX, LY, LW = PAD_X, 2.7, 4.5
    add_text(s, LX, LY, LW, 0.3,
             "저평점 원인 분포 (4팀 분산)",
             size=10, font=FONT_MEDIUM, color=GRAY_600, char_spacing=1.4)
    BAR_Y = LY + 0.45
    BAR_H = 0.6
    segs = [(42, OLIST_BLUE, "물류"), (20, GRAY_800, "CX"),
            (20, GRAY_600, "SO"), (18, GRAY_400, "Prod")]
    sx = LX
    for pct, color, label in segs:
        sw = LW * (pct / 100.0)
        add_rect(s, sx, BAR_Y, sw, BAR_H, fill=color)
        add_text(s, sx, BAR_Y + 0.1, sw, 0.3, f"{pct}%",
                 size=14, font=FONT_BLACK, bold=True, color=WHITE, align="center")
        add_text(s, sx, BAR_Y + BAR_H + 0.05, sw, 0.25, label,
                 size=9, font=FONT_REGULAR, color=GRAY_600, align="center")
        sx += sw
    # Caption
    add_text(s, LX, LY + 1.6, LW, 0.4,
             "원인 100%가 4개 팀에 분산.\n한 팀 단독으로 풀 수 없는 분포.",
             size=11, font=FONT_REGULAR, color=GRAY_600, line_spacing=1.5)

    # 우측 TF 매핑 3카드
    RX, RY, RW = 7.5, 2.7, 5.0
    add_text(s, RX, RY, RW, 0.3,
             "TF 매핑 (3-Core)",
             size=10, font=FONT_MEDIUM, color=GRAY_600, char_spacing=1.4)
    tfs = [
        {"id": "TF1", "title": "배송 품질", "teams": "Logistics · Seller Ops · CX",
         "y": RY + 0.45, "highlight": False},
        {"id": "TF2", "title": "구매 경험", "teams": "Product · 물류 · CX · SO",
         "y": RY + 1.45, "highlight": False},
        {"id": "TF3 ★", "title": "재구매 유도", "teams": "CRM · Product · MD",
         "y": RY + 2.45, "highlight": True},
    ]
    for t in tfs:
        bg = BLUE_TINT if t["highlight"] else WHITE
        line = OLIST_BLUE if t["highlight"] else GRAY_200
        add_rect(s, RX, t["y"], RW, 0.85, fill=bg, line=line, line_width=1)
        if t["highlight"]:
            add_rect(s, RX, t["y"], 0.05, 0.85, fill=OLIST_BLUE)
        add_text(s, RX + 0.2, t["y"] + 0.1, 1.0, 0.4, t["id"],
                 size=18, font=FONT_BLACK, bold=True,
                 color=OLIST_BLUE if t["highlight"] else GRAY_600)
        add_text(s, RX + 1.4, t["y"] + 0.1, RW - 1.6, 0.4, t["title"],
                 size=15, font=FONT_BOLD, bold=True, color=NEAR_BLACK)
        add_text(s, RX + 1.4, t["y"] + 0.5, RW - 1.6, 0.3, t["teams"],
                 size=10, font=FONT_REGULAR, color=GRAY_600)

    # 하단 Highlight
    add_rect(s, PAD_X, 6.55, 12.21, 0.55, fill=BLUE_TINT)
    add_rect(s, PAD_X, 6.55, 0.05, 0.55, fill=OLIST_BLUE)
    add_text(s, PAD_X + 0.3, 6.7, 12, 0.3,
             "TF3 ★는 발표 핵심 KR(KR1 재구매율 4.5%)을 직접 책임지는 TF입니다.",
             size=12, font=FONT_BOLD, bold=True, color=NEAR_BLACK)

    add_footer(s, page_num=11)
    return s


# ============================================================
# Slide 12 — Type 8: 5팀 × 3TF 매트릭스
# ============================================================
def make_slide_12(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_header(s, active_idx=2)

    add_page_title(s, runs=[("FIVE TEAMS  ", NEAR_BLACK),
                            ("×", OLIST_BLUE),
                            ("  THREE TFs", NEAR_BLACK)], size=42)
    add_lead(s, "3-Core 팀이 각 TF 리드, Support 2팀(SO·MD)이 교차 보조.")

    # 매트릭스 표 (ROW_H 0.7 → 0.55로 줄여 5행 + Highlight Box 안 겹침)
    HX, HY = PAD_X + 2.5, 2.7
    COL_W = 3.0
    ROW_H = 0.55
    teams = ["CRM", "Product/UX", "Logistics", "Seller Ops", "MD"]
    tf_ids = ["TF1 배송 품질", "TF2 구매 경험", "TF3 재구매 유도 ★"]
    matrix = [
        ["–",   "–",   "🔵 리드 ★"],   # CRM
        ["보조", "🔵 리드 ★", "보조"],   # Product
        ["🔵 리드 ★", "보조", "–"],     # Logistics
        ["보조", "보조", "–"],          # Seller Ops
        ["–",   "–",   "보조"],         # MD
    ]
    # 헤더 (TF)
    for j, tf in enumerate(tf_ids):
        bx = HX + j * COL_W
        is_tf3 = (j == 2)
        if is_tf3:
            add_rect(s, bx, HY, COL_W, ROW_H, fill=BLUE_TINT)
        else:
            add_rect(s, bx, HY, COL_W, ROW_H, fill=GRAY_200)
        add_text(s, bx + 0.2, HY + 0.13, COL_W - 0.4, 0.3, tf,
                 size=12, font=FONT_BOLD, bold=True,
                 color=OLIST_BLUE if is_tf3 else NEAR_BLACK, align="center")
    # 행 (팀)
    for i, team in enumerate(teams):
        ry = HY + (i + 1) * ROW_H + 0.05
        # Team name (좌)
        add_text(s, PAD_X, ry + 0.13, 2.4, 0.3, team,
                 size=12, font=FONT_BOLD, bold=True, color=NEAR_BLACK)
        # Cell
        for j, val in enumerate(matrix[i]):
            cx = HX + j * COL_W
            is_lead = "리드" in val
            is_tf3_lead = is_lead and j == 2
            if is_tf3_lead:
                add_rect(s, cx, ry, COL_W, ROW_H, fill=BLUE_TINT)
            elif is_lead:
                add_rect(s, cx, ry, COL_W, ROW_H, fill=GRAY_200)
            color = OLIST_BLUE if is_tf3_lead else (NEAR_BLACK if is_lead else GRAY_400)
            add_text(s, cx + 0.2, ry + 0.13, COL_W - 0.4, 0.3, val,
                     size=11,
                     font=FONT_BOLD if is_lead else FONT_REGULAR,
                     bold=is_lead, color=color, align="center")
        # 가로 구분선
        add_h_line(s, PAD_X, ry + ROW_H + 0.02, 12.21)

    # 하단 Highlight
    add_rect(s, PAD_X, 6.55, 12.21, 0.55, fill=BLUE_TINT)
    add_rect(s, PAD_X, 6.55, 0.05, 0.55, fill=OLIST_BLUE)
    add_text(s, PAD_X + 0.3, 6.7, 12, 0.3,
             "★ 표시 = 발표 정점. TF3 컬럼이 핵심 KR을 직접 책임지는 영역.",
             size=12, font=FONT_BOLD, bold=True, color=NEAR_BLACK)

    add_footer(s, page_num=12)
    return s


# ============================================================
# Slide 13 — Type 6: 이중 트랙 (Dual Track)
# ============================================================
def make_slide_13(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_header(s, active_idx=2)

    add_page_title(s, runs=[("DUAL TRACK  ", NEAR_BLACK),
                            ("TO 4.5%", OLIST_BLUE)], size=44)
    add_lead(s, "TF는 BF로 증명, 구조 개선은 1년으로 완성. 둘 다 KR1 4.5%로 수렴.")

    # 단기 트랙 (위)
    SX = PAD_X + 1.8
    SW = 9.5
    SHORT_Y = 3.2
    add_text(s, PAD_X, SHORT_Y, 1.6, 0.3,
             "단기 트랙 — TF",
             size=10, font=FONT_MEDIUM, color=OLIST_BLUE, char_spacing=1.4)
    add_h_line(s, SX, SHORT_Y + 0.2, SW, color=OLIST_BLUE)
    s_nodes = [
        ("01", "Now", "TF 셋업"),
        ("02", "BF (6M)", "cohort 1.5~2.0% ★"),
        ("03", "12M", "내재화"),
    ]
    for i, (n, p, d) in enumerate(s_nodes):
        cx = SX + i * (SW / (len(s_nodes) - 1))
        add_circle(s, cx, SHORT_Y + 0.2, 0.4, fill=OLIST_BLUE)
        add_text(s, cx - 0.2, SHORT_Y + 0.1, 0.4, 0.25, n,
                 size=10, font=FONT_BOLD, bold=True, color=WHITE, align="center")
        add_text(s, cx - 1.0, SHORT_Y - 0.45, 2.0, 0.3, p,
                 size=11, font=FONT_BOLD, bold=True, color=NEAR_BLACK, align="center")
        add_text(s, cx - 1.0, SHORT_Y + 0.55, 2.0, 0.3, d,
                 size=10, font=FONT_REGULAR, color=GRAY_600, align="center")

    # 장기 트랙 (아래)
    LONG_Y = 4.8
    add_text(s, PAD_X, LONG_Y, 1.6, 0.3,
             "장기 트랙 — 구조 개선",
             size=10, font=FONT_MEDIUM, color=GRAY_600, char_spacing=1.4)
    add_h_line(s, SX, LONG_Y + 0.2, SW, color=GRAY_400)
    l_nodes = [
        ("a", "물류 3PL", "지역 확장"),
        ("b", "SO 등급제", "C → A 셀러"),
        ("c", "MD 번들", "카테고리 다각화"),
    ]
    for i, (n, p, d) in enumerate(l_nodes):
        cx = SX + i * (SW / (len(l_nodes) - 1))
        add_circle(s, cx, LONG_Y + 0.2, 0.4, fill=WHITE, line=GRAY_400)
        add_text(s, cx - 0.2, LONG_Y + 0.1, 0.4, 0.25, n,
                 size=10, font=FONT_BOLD, bold=True, color=GRAY_600, align="center")
        add_text(s, cx - 1.0, LONG_Y - 0.45, 2.0, 0.3, p,
                 size=11, font=FONT_BOLD, bold=True, color=NEAR_BLACK, align="center")
        add_text(s, cx - 1.0, LONG_Y + 0.55, 2.0, 0.3, d,
                 size=10, font=FONT_REGULAR, color=GRAY_600, align="center")

    # 우측 KR1 4.5% 수렴
    KX = SX + SW + 0.3
    KY = 3.6
    add_text(s, KX, KY, 1.5, 0.3,
             "→ 수렴",
             size=10, font=FONT_MEDIUM, color=OLIST_BLUE, char_spacing=1.4)
    add_text(s, KX, KY + 0.4, 1.5, 0.6,
             "KR1",
             size=20, font=FONT_BOLD, bold=True, color=NEAR_BLACK)
    add_text(s, KX, KY + 0.95, 1.5, 0.7,
             "4.5%",
             size=44, font=FONT_BLACK, bold=True, color=OLIST_BLUE)

    # 하단 Highlight
    add_rect(s, PAD_X, 6.55, 12.21, 0.55, fill=BLUE_TINT)
    add_rect(s, PAD_X, 6.55, 0.05, 0.55, fill=OLIST_BLUE)
    add_text(s, PAD_X + 0.3, 6.7, 12, 0.3,
             "두 트랙 모두 KR1 4.5%로 수렴 — 그 수렴점을 직접 책임지는 게 TF3.",
             size=12, font=FONT_BOLD, bold=True, color=NEAR_BLACK)

    add_footer(s, page_num=13)
    return s


# ============================================================
# Slide 14 — Type 4: TF1 ① 지표·목표
# ============================================================
def make_slide_14(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_header(s, active_idx=2)

    add_page_title(s, runs=[("TF1  ", NEAR_BLACK),
                            ("DELIVERY", OLIST_BLUE)], w=7.0, size=44)
    # Meta (우측, Page Title 끝점 7.56" 와 갭 확보)
    add_text(s, RIGHT_EDGE - 5, 1.20, 5, 0.3,
             "AARRR — ACTIVATION 🚨  ·  KR4 책임 (간접)",
             size=9, font=FONT_MEDIUM, color=GRAY_600,
             char_spacing=1.4, align="right")
    add_lead(s, "TF1은 배송 품질. 첫 구매 지연율을 6.0%로, BF 지연율을 15% 이하로 통제.")

    # 2 KPI 게이지
    add_kpi_gauge(s, PAD_X, 2.7, 5.85,
                  "KR4 — 첫 구매 배송 지연율", "8.16%", "6.0%", progress=0.45, big=True)
    add_kpi_gauge(s, PAD_X + 6.36, 2.7, 5.85,
                  "KR2 — BF 기간 전체 지연율", "20.93%", "15.0% 이하", progress=0.40, big=True)

    # 마일스톤
    add_text(s, PAD_X, 4.6, 12.21, 0.3,
             "MILESTONE  —  3M / 6M / 12M",
             size=10, font=FONT_MEDIUM, color=GRAY_600, char_spacing=1.4)
    add_h_line(s, PAD_X, 4.95, 12.21)
    add_milestone_strip(s, PAD_X, 5.1, 12.21, items=[
        ("3M", "지연 다발 셀러 모니터링", "Top 100 셀러 SLA 감시 + 자동 알람"),
        ("6M", "북동부 3PL 파트너십", "RJ·MA 등 매출 高 지역 1P 보강"),
        ("12M", "Lead Time 단축", "Total Lead Time 12.1일 → 8일"),
    ])

    add_footer(s, page_num=14)
    return s


# ============================================================
# Slide 15 — Type 5: TF1 Before/After
# ============================================================
def make_slide_15(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_header(s, active_idx=2)

    add_page_title(s, runs=[("TF1  ", NEAR_BLACK),
                            ("BEFORE / AFTER", OLIST_BLUE)], size=44)
    add_lead(s, "병목 — BF 7.94배 폭증으로 지연율 20.93% 도달. TF1 가동으로 평시·BF 모두 통제.")

    # 2열 카드
    CW = (12.21 - 0.4) / 2
    add_comparison_card(s, PAD_X, 2.8, CW, 3.7,
        label="BEFORE — 평시 / BF",
        head="지연율\n8.16% / 20.93%",
        body="• BF 주문량 7.94배 폭증 시 지연 통제 불가\n"
             "• 첫 구매 고객의 부정 리뷰 확산\n"
             "• 매출 高 지역(북동부)에 셀러 부족",
        stat_label="평시 / BF", stat_value="8.16 / 20.93", highlight=False)
    add_comparison_card(s, PAD_X + CW + 0.4, 2.8, CW, 3.7,
        label="AFTER — TF1 가동",
        head="지연율\n6.0% / 15.0% 이하",
        body="• 지연 다발 셀러 사전 모니터링·교육\n"
             "• 북동부 3PL 파트너십 (RJ·MA)\n"
             "• Total Lead Time 12.1일 → 8일",
        stat_label="평시 / BF", stat_value="6.0 / 15.0", highlight=True)

    # 하단 화살표 안내
    add_text(s, PAD_X, 6.7, 12.21, 0.3,
             "→ BF 지연율 -5.93%p 감축 = 첫 구매 고객 만족도 회복 + TF2 첫 경험 KR로 연결",
             size=11, font=FONT_REGULAR, color=GRAY_600, align="center")

    add_footer(s, page_num=15)
    return s


# ============================================================
# Slide 16 — Type 4: TF1 결과물 (지도 4종 + 번들 3종 + 셀러 교육)
# ============================================================
def make_slide_16(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_header(s, active_idx=2)

    add_page_title(s, runs=[("TF1  ", NEAR_BLACK),
                            ("OUTPUTS", OLIST_BLUE)], size=44)
    add_lead(s, "지역별 지연율 지도 4종 + 셀러 교육 압축안 + MD 묶음 배송 번들 3종.")

    # 좌측: 브라질 지도 4장 (2×2 미니 모자이크)
    LX, LY, LW, LH = PAD_X, 2.7, 5.5, 4.1
    add_text(s, LX, LY, LW, 0.3,
             "OUTPUT 1 — 브라질 지도 (지연·셀러·Gap·1P)",
             size=10, font=FONT_MEDIUM, color=GRAY_600, char_spacing=1.4)
    map_paths = [
        "발표자료/TF1-배송품질/assets/tf1-brazilmap1-delay.png",
        "발표자료/TF1-배송품질/assets/tf1-brazilmap2-sellers count.png",
        "발표자료/TF1-배송품질/assets/tf1-brazilmap3-gap(supply_demand).png",
        "발표자료/TF1-배송품질/assets/tf1-brazilmap4-1p-spots.png",
    ]
    map_w = (LW - 0.15) / 2
    map_h = (LH - 0.55) / 2
    for i, p in enumerate(map_paths):
        mx = LX + (i % 2) * (map_w + 0.15)
        my = LY + 0.4 + (i // 2) * (map_h + 0.15)
        add_picture_safe(s, p, mx, my, map_w, map_h)

    # 우측 상단: 셀러 교육 4 대상
    EX, EY, EW = 6.3, 2.7, 6.21
    EH = 1.95
    add_text(s, EX, EY, EW, 0.3,
             "OUTPUT 2 — 셀러 교육 (4 대상 압축)",
             size=10, font=FONT_MEDIUM, color=GRAY_600, char_spacing=1.4)
    edu = [
        ("A", "신규 셀러", "첫 30일 가이드"),
        ("B", "기존 셀러", "지연 vs 리뷰"),
        ("C", "위험 셀러", "1점 대응법"),
        ("D", "VIP 셀러", "전담 AM"),
    ]
    cw = (EW - 0.3) / 4
    for i, (g, t, d) in enumerate(edu):
        cx = EX + i * (cw + 0.1)
        add_rect(s, cx, EY + 0.4, cw, EH - 0.4, fill=GRAY_200)
        add_text(s, cx + 0.15, EY + 0.55, cw - 0.3, 0.4, g,
                 size=20, font=FONT_BLACK, bold=True, color=GRAY_600)
        add_text(s, cx + 0.15, EY + 1.0, cw - 0.3, 0.3, t,
                 size=11, font=FONT_BOLD, bold=True, color=NEAR_BLACK)
        add_text(s, cx + 0.15, EY + 1.3, cw - 0.3, 0.3, d,
                 size=9, font=FONT_REGULAR, color=GRAY_600)

    # 우측 하단: MD 번들 3종
    BX, BY, BW = 6.3, 4.85, 6.21
    BH = 1.85
    add_text(s, BX, BY, BW, 0.3,
             "OUTPUT 3 — MD 묶음 배송 번들 (BF 배송비 흡수)",
             size=10, font=FONT_MEDIUM, color=GRAY_600, char_spacing=1.4)
    bundle_paths = [
        "발표자료/TF1-배송품질/assets/bundle_banner.png",
        "발표자료/TF1-배송품질/assets/bundle_proposal_visual.png",
        "발표자료/TF1-배송품질/assets/bundle_sdp.png",
    ]
    bw = (BW - 0.3) / 3
    bh = BH - 0.4
    for i, p in enumerate(bundle_paths):
        bx = BX + i * (bw + 0.15)
        add_picture_safe(s, p, bx, BY + 0.4, bw, bh)

    add_footer(s, page_num=16)
    return s


# ============================================================
# Slide 17 — Type 4: TF2 ① 지표·목표
# ============================================================
def make_slide_17(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_header(s, active_idx=2)

    add_page_title(s, runs=[("TF2  ", NEAR_BLACK),
                            ("EXPERIENCE", OLIST_BLUE)], w=7.0, size=44)
    add_text(s, RIGHT_EDGE - 5, 1.20, 5, 0.3,
             "AARRR — ACTIVATION ⚠  ·  KR5 (간접) · 4팀 공동",
             size=9, font=FONT_MEDIUM, color=GRAY_600,
             char_spacing=1.4, align="right")
    add_lead(s, "TF2는 구매 경험. 저평점 12.82% → 10.5%, 4팀이 공동 KR 1개 + 기여 KR 4개로 묶어 푼다.")

    # 1 메인 KPI + 4 기여 KR
    add_kpi_gauge(s, PAD_X, 2.7, 12.21,
                  "공동 KR — 첫 구매 저평점 비중", "12.82%", "10.5%",
                  progress=0.55, big=True)

    # 4 기여 KR 카드
    sub_y = 4.4
    sub_h = 1.6
    sub_w = (12.21 - 3 * 0.2) / 4
    subs = [
        {"team": "물류", "label": "KR4", "head": "8.16% → 6.0%", "desc": "지연 통제"},
        {"team": "CX", "label": "응답 SLA", "head": "48h 90%", "desc": "대응률 달성"},
        {"team": "Seller Ops", "label": "위험 셀러", "head": "23 → 15명", "desc": "4.0점 회복"},
        {"team": "Product/UX", "label": "PDP·결제", "head": "신뢰 UX", "desc": "기대 일치"},
    ]
    for i, sub in enumerate(subs):
        cx = PAD_X + i * (sub_w + 0.2)
        add_rect(s, cx, sub_y, sub_w, sub_h, fill=GRAY_200)
        add_text(s, cx + 0.2, sub_y + 0.15, sub_w - 0.4, 0.25, sub["team"],
                 size=9, font=FONT_MEDIUM, color=GRAY_600, char_spacing=1.4)
        add_text(s, cx + 0.2, sub_y + 0.45, sub_w - 0.4, 0.3, sub["label"],
                 size=11, font=FONT_BOLD, bold=True, color=NEAR_BLACK)
        add_text(s, cx + 0.2, sub_y + 0.8, sub_w - 0.4, 0.4, sub["head"],
                 size=18, font=FONT_BLACK, bold=True, color=OLIST_BLUE)
        add_text(s, cx + 0.2, sub_y + 1.25, sub_w - 0.4, 0.3, sub["desc"],
                 size=10, font=FONT_REGULAR, color=GRAY_600)

    # 마일스톤
    add_text(s, PAD_X, 6.2, 12.21, 0.3,
             "MILESTONE  —  3M PDP / 6M Thank-you / 12M 상품 정보 표준화",
             size=10, font=FONT_MEDIUM, color=GRAY_600, char_spacing=1.4)

    add_footer(s, page_num=17)
    return s


# ============================================================
# Slide 18 — Type 5: TF2 Before/After
# ============================================================
def make_slide_18(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_header(s, active_idx=2)

    add_page_title(s, runs=[("TF2  ", NEAR_BLACK),
                            ("BEFORE / AFTER", OLIST_BLUE)], size=44)
    add_lead(s, "병목 — 저평점 원인이 4팀에 분산. 공동 KR 1개 + 기여 KR 4개로 책임을 명확히.")

    CW = (12.21 - 0.4) / 2
    add_comparison_card(s, PAD_X, 2.8, CW, 3.7,
        label="BEFORE",
        head="책임 모호\n저평점 12.82%",
        body="• 원인 4팀 분산 → 단독 풀이 불가\n"
             "• 책임 소재 모호 (누가 KR을 만들지?)\n"
             "• 평점·재구매 핸드오프 없음",
        stat_label="저평점 비중", stat_value="12.82%", highlight=False)
    add_comparison_card(s, PAD_X + CW + 0.4, 2.8, CW, 3.7,
        label="AFTER — TF2 가동",
        head="공동 KR + 기여 KR\n저평점 10.5%",
        body="• 4팀 공동 KR 1개 + 기여 KR 4개\n"
             "• 월 1회 합동 리뷰 (책임 가시화)\n"
             "• 평점 ↑ → 재구매 (TF3 핸드오프)",
        stat_label="저평점 비중", stat_value="10.5%", highlight=True)

    add_text(s, PAD_X, 6.7, 12.21, 0.3,
             "→ 메커니즘: 유저 만족 ↑ → 평점 ↑ → 5점 재구매율 3.12%로 연결 (TF3 진입점)",
             size=11, font=FONT_REGULAR, color=GRAY_600, align="center")

    add_footer(s, page_num=18)
    return s


# ============================================================
# Slide 19 — Type 4: TF2 결과물 4종
# ============================================================
def make_slide_19(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_header(s, active_idx=2)

    add_page_title(s, runs=[("TF2  ", NEAR_BLACK),
                            ("OUTPUTS", OLIST_BLUE)], size=44)
    add_lead(s, "PDP 신뢰 + CX 48h SLA + Thank-you 4채널 + BF 랜딩 — 4종 산출물.")

    # 2×2 그리드 + Thank-you 영역만 4 mini
    GX, GY = PAD_X, 2.7
    cw = (12.21 - 0.3) / 2
    ch = (4.3 - 0.2) / 2

    # 좌상 — PDP After 메인
    add_text(s, GX, GY, cw, 0.3, "OUTPUT 1 — PDP After 메인",
             size=9, font=FONT_MEDIUM, color=GRAY_600, char_spacing=1.4)
    add_picture_safe(s, "발표자료/TF2-구매경험/assets/tf2-pdp-hero-soft-curve.png",
                     GX, GY + 0.35, cw, ch - 0.4)

    # 우상 — CX SLA
    add_text(s, GX + cw + 0.3, GY, cw, 0.3, "OUTPUT 2 — CX 48h SLA (90%)",
             size=9, font=FONT_MEDIUM, color=GRAY_600, char_spacing=1.4)
    add_picture_safe(s, "발표자료/TF2-구매경험/assets/tf2-cx-sla-90percent.png",
                     GX + cw + 0.3, GY + 0.35, cw, ch - 0.4)

    # 좌하 — Thank-you 4채널 mini grid (D+0 → D+60)
    THX = GX
    THY = GY + ch + 0.2
    add_text(s, THX, THY, cw, 0.3,
             "OUTPUT 3 — Thank-you 4채널 (D+0 → D+60)",
             size=9, font=FONT_MEDIUM, color=GRAY_600, char_spacing=1.4)
    thankyou = [
        ("D+0 After", "발표자료/TF2-구매경험/assets/tf2-thankyou-after.png"),
        ("D+7 Push", "발표자료/TF2-구매경험/assets/tf2-thankyou-push-d7.png"),
        ("D+30 Email", "발표자료/TF2-구매경험/assets/tf2-thankyou-email-d30.png"),
        ("D+60 Kakao", "발표자료/TF2-구매경험/assets/tf2-thankyou-kakao-d60.png"),
    ]
    mini_w = (cw - 0.3) / 4
    mini_label_h = 0.2
    mini_img_h = ch - 0.4 - mini_label_h
    for i, (label, img) in enumerate(thankyou):
        mx = THX + i * (mini_w + 0.1)
        add_text(s, mx, THY + 0.35, mini_w, mini_label_h, label,
                 size=8, font=FONT_BOLD, bold=True, color=OLIST_BLUE,
                 char_spacing=0.6, align="center")
        add_picture_safe(s, img, mx, THY + 0.35 + mini_label_h, mini_w, mini_img_h)

    # 우하 — BF 랜딩 ref
    add_text(s, GX + cw + 0.3, THY, cw, 0.3, "OUTPUT 4 — BF 랜딩 ref",
             size=9, font=FONT_MEDIUM, color=GRAY_600, char_spacing=1.4)
    add_picture_safe(s, "발표자료/TF2-구매경험/assets/tf2-bf-ref-gmarket-blackprime.png",
                     GX + cw + 0.3, THY + 0.35, cw, ch - 0.4)

    add_footer(s, page_num=19)
    return s


# ============================================================
# Slide 20 — Type 4: TF3 ① 지표·목표 ★
# ============================================================
def make_slide_20(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_header(s, active_idx=2)

    add_page_title(s, runs=[("TF3  ", NEAR_BLACK),
                            ("RETENTION ★", OLIST_BLUE)], w=7.0, size=44)
    add_text(s, RIGHT_EDGE - 5.5, 1.20, 5.5, 0.3,
             "AARRR — RETENTION 🚨  ·  KR1·KR3 직접 책임 ★",
             size=9, font=FONT_MEDIUM, color=GRAY_600,
             char_spacing=1.4, align="right")
    add_lead(s, "TF3는 재구매 유도. KR1·KR3·BF cohort 3개를 직접 책임 — 비중 2분.")

    # KR1 메인 (큰 게이지) + KR3·BF 서브
    add_rect(s, PAD_X, 2.7, 12.21, 1.7, fill=BLUE_TINT)
    add_rect(s, PAD_X, 2.7, 0.05, 1.7, fill=OLIST_BLUE)
    add_text(s, PAD_X + 0.3, 2.85, 11.91, 0.3,
             "KR1 ★ — 재구매율 (전사 핵심 KR)",
             size=11, font=FONT_MEDIUM, color=OLIST_BLUE, char_spacing=1.4)
    add_text(s, PAD_X + 0.3, 3.2, 5, 0.7, "현재 3.0%",
             size=36, font=FONT_BOLD, bold=True, color=NEAR_BLACK)
    # 목표 4.5% — Semantic SUCCESS 적용
    add_text(s, PAD_X + 6, 3.2, 6, 0.7, "목표 4.5% ✓",
             size=36, font=FONT_BLACK, bold=True, color=SUCCESS, align="right")
    add_rect(s, PAD_X + 0.3, 4.05, 11.6, 0.08, fill=GRAY_200)
    add_rect(s, PAD_X + 0.3, 4.05, 11.6 * 0.4, 0.08, fill=OLIST_BLUE)

    # 2 sub KR
    sub_y = 4.7
    sub_h = 1.4
    sub_w = (12.21 - 0.3) / 2
    subs = [
        ("KR3 — 재구매 주문 비중", "6.13%", "9.0%", 0.45),
        ("BF-KR5 — BF cohort 30일 재구매", "0.56%", "1.5~2.0%", 0.30),
    ]
    for i, (lbl, c, t, p) in enumerate(subs):
        cx = PAD_X + i * (sub_w + 0.3)
        add_kpi_gauge(s, cx, sub_y, sub_w, lbl, c, t, progress=p, big=False)

    # 마일스톤
    add_text(s, PAD_X, 6.45, 12.21, 0.3,
             "MILESTONE  —  3M D+7 푸시 / 6M D+30 자동화 / 12M 7터치포인트 자동화",
             size=10, font=FONT_MEDIUM, color=GRAY_600, char_spacing=1.4)

    add_footer(s, page_num=20)
    return s


# ============================================================
# Slide 21 — Type 5: TF3 Before/After
# ============================================================
def make_slide_21(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_header(s, active_idx=2)

    add_page_title(s, runs=[("TF3  ", NEAR_BLACK),
                            ("BEFORE / AFTER", OLIST_BLUE)], size=44)
    add_lead(s, "병목 — 재구매 유도 인프라 전무 (CRM·추천·쿠폰 모두 없음).")

    CW = (12.21 - 0.4) / 2
    add_comparison_card(s, PAD_X, 2.8, CW, 3.7,
        label="BEFORE",
        head="공백 → 97% 이탈",
        body="• 구매 완료 → 어떤 훅도 없음\n"
             "• CRM · 추천 · 쿠폰 인프라 전무\n"
             "• 재방문 의향 만들 도구 부재",
        stat_label="재구매율", stat_value="3.0%", highlight=False)
    add_comparison_card(s, PAD_X + CW + 0.4, 2.8, CW, 3.7,
        label="AFTER — TF3 가동",
        head="7터치포인트 자동화",
        body="• D+0 ~ D+90 라이프사이클 자동화\n"
             "• CRM + MD + Product 협업\n"
             "• 코호트 기반 재방문 트리거 + 쿠폰",
        stat_label="재구매율 (1년)", stat_value="4.5%", highlight=True)

    # 정점 멘트
    add_rect(s, PAD_X, 6.55, 12.21, 0.55, fill=BLUE_TINT)
    add_rect(s, PAD_X, 6.55, 0.05, 0.55, fill=OLIST_BLUE)
    add_text(s, PAD_X + 0.3, 6.7, 12, 0.3,
             "★ 앞 모든 분석이 결국 한 가지로 수렴 — 재구매 인프라가 아예 없다는 사실.",
             size=12, font=FONT_BOLD, bold=True, color=OLIST_BLUE)

    add_footer(s, page_num=21)
    return s


# ============================================================
# Slide 23 — Type 4: TF3 결과물 4종 + CRM best-of 3장
# ============================================================
def make_slide_23(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_header(s, active_idx=2)

    add_page_title(s, runs=[("TF3  ", NEAR_BLACK),
                            ("OUTPUTS", OLIST_BLUE)], size=44)
    add_lead(s, "CRM 목업 best-of 3장 + 번들 + 셀러 다각화 + BF 코호트 — 7터치포인트의 실체.")

    # 상단: CRM 3장
    TX, TY, TW = PAD_X, 2.7, 12.21
    add_text(s, TX, TY, TW, 0.3,
             "CRM Best-of 3 — D+7 / D+30 / D+32",
             size=10, font=FONT_MEDIUM, color=GRAY_600, char_spacing=1.4)
    crm = [
        ("발표자료/TF3-재구매/assets/tf3_D07_push_A.png", "D+7 푸시"),
        ("발표자료/TF3-재구매/assets/tf3_D30_kakao.png", "D+30 카카오"),
        ("발표자료/TF3-재구매/assets/tf3_D32_kakao.png", "D+32 카카오"),
    ]
    img_w = (TW - 0.4) / 3
    img_h = 2.05
    for i, (p, name) in enumerate(crm):
        ix = TX + i * (img_w + 0.2)
        add_text(s, ix, TY + 0.4, img_w, 0.25, name,
                 size=9, font=FONT_BOLD, bold=True, color=NEAR_BLACK)
        add_picture_safe(s, p, ix, TY + 0.7, img_w, img_h)

    # 하단: 4 영역 (3 이미지 + 1 캡션)
    BY = 5.45
    BH = 1.6
    bw = (TW - 3 * 0.18) / 4
    bh = BH - 0.4
    add_text(s, TX, BY, TW, 0.3,
             "산출물 4종",
             size=10, font=FONT_MEDIUM, color=GRAY_600, char_spacing=1.4)
    bottom = [
        {"img": "발표자료/TF3-재구매/assets/tf3-bundle-bedding-set.png", "label": "OUT 2 — 번들 카드"},
        {"img": "발표자료/TF3-재구매/assets/tf3-seller-diversification.png", "label": "OUT 3 — 셀러 다각화"},
        {"img": "발표자료/TF3-재구매/assets/tf3-bf-cohort-tracking.png", "label": "OUT 4 — BF 코호트"},
        {"img": None, "label": "★ 7터치의 실체"},
    ]
    for i, it in enumerate(bottom):
        ix = TX + i * (bw + 0.18)
        add_text(s, ix, BY + 0.35, bw, 0.22, it["label"],
                 size=8, font=FONT_MEDIUM, color=GRAY_600, char_spacing=1.0)
        if it["img"]:
            add_picture_safe(s, it["img"], ix, BY + 0.6, bw, bh - 0.25)
        else:
            add_rect(s, ix, BY + 0.6, bw, bh - 0.25, fill=BLUE_TINT)
            add_rect(s, ix, BY + 0.6, 0.04, bh - 0.25, fill=OLIST_BLUE)
            add_text(s, ix + 0.15, BY + 0.7, bw - 0.3, 0.4,
                     "CRM·MD·Product가\n한 화면에",
                     size=12, font=FONT_BOLD, bold=True,
                     color=NEAR_BLACK, line_spacing=1.3)
            add_text(s, ix + 0.15, BY + 1.15, bw - 0.3, 0.3,
                     "모이는 지점",
                     size=11, font=FONT_REGULAR, color=GRAY_600)

    add_footer(s, page_num=23)
    return s


# ============================================================
# Slide 24 — Type 5: 결론 재선언
# ============================================================
def make_slide_24(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_header(s, active_idx=3)  # Outcome

    add_page_title(s, runs=[("FROM  ", NEAR_BLACK),
                            ("NO STRUCTURE", OLIST_BLUE),
                            ("  TO 4.5%", NEAR_BLACK)], size=42)
    add_lead(s, "결론 재선언 — 슬4와 같은 카드, 다른 답.")

    CW = (12.21 - 0.4) / 2
    add_comparison_card(s, PAD_X, 2.8, CW, 3.7,
        label="BEFORE — 슬4의 결론",
        head="재구매 구조가\n없습니다",
        body="• 재구매율 3.0%\n• BF cohort 30일 0.56%\n"
             "• 재구매 주문 비중 6.13%\n• 매출 94%가 일회성 의존",
        stat_label="현재 KR1", stat_value="3.0%", highlight=False)
    add_comparison_card(s, PAD_X + CW + 0.4, 2.8, CW, 3.7,
        label="AFTER — TF 3개로",
        head="재구매 구조를\n만듭니다",
        body="• TF1 배송 (지연 8.16% → 6.0%)\n"
             "• TF2 첫 경험 (저평점 12.82% → 10.5%)\n"
             "• TF3 ★ 재구매 유도 (KR1 직접)\n"
             "• BF 1.5~2.0% → 12M 4.5% 수렴",
        stat_label="1년 후 KR1", stat_value="4.5%", highlight=True)

    add_text(s, PAD_X, 6.7, 12.21, 0.3,
             "→ 슬4에서 드린 결론을 다시 한 번 — 우리는 이렇게 풀겠습니다.",
             size=12, font=FONT_BOLD, bold=True, color=NEAR_BLACK, align="center")

    add_footer(s, page_num=24)
    return s


# ============================================================
# Slide 25 — Type 8: 기대효과 표 (3M / 6M / 12M)
# ============================================================
def make_slide_25(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_header(s, active_idx=3)

    add_page_title(s, runs=[("EXPECTED  ", NEAR_BLACK),
                            ("IMPACT", OLIST_BLUE)], w=8, size=44)
    add_text(s, RIGHT_EDGE - 5, 1.20, 5, 0.3,
             "3M / 6M / 12M  ·  By 책임 TF",
             size=9, font=FONT_MEDIUM, color=GRAY_600,
             char_spacing=1.4, align="right")
    add_lead(s, "3개월 = TF 성과, 6~12개월 = 장기 구조 개선 성과.")

    # 표
    HX, HY = PAD_X, 2.7
    cols_w = [1.4, 3.6, 1.35, 1.5, 1.5, 1.5]  # 합 10.85, padding 포함
    # 사용 영역: PAD_X + 11 = 11.56 (RIGHT_EDGE 12.5 안쪽)
    headers = ["TF", "INDICATOR", "NOW", "3M (BF)", "6M", "12M"]
    cx = HX
    for w, h in zip(cols_w, headers):
        add_text(s, cx + 0.1, HY, w - 0.1, 0.3, h,
                 size=9, font=FONT_MEDIUM, color=GRAY_600, char_spacing=1.4)
        cx += w
    add_h_line(s, HX, HY + 0.4, sum(cols_w))

    rows = [
        ("★ TF3", "KR1 재구매율 (전사 핵심)", "3.0%", "3.5%", "4.0%", "4.5%", True),
        ("★ TF3", "BF 30일 재구매", "0.56%", "1.5~2.0%", "–", "–", True),
        ("★ TF3", "KR3 재구매 주문 비중", "6.13%", "7.0%", "8.0%", "9.0%", True),
        ("TF1", "첫 구매 지연율", "8.16%", "7.0%", "6.5%", "6.0%", False),
        ("TF1 (장기)", "Total Lead Time", "12.1일", "–", "10일", "8일", False),
        ("TF2", "첫 구매 저평점", "12.82%", "11.8%", "11.0%", "10.5%", False),
        ("TF1 (셀러)", "위험 셀러 4.0점", "0/23", "5/23", "15/23", "–", False),
    ]
    ROW_H = 0.45
    ry = HY + 0.5
    for r in rows:
        if r[6]:
            add_rect(s, HX, ry, sum(cols_w), ROW_H, fill=BLUE_TINT)
            add_rect(s, HX, ry, 0.05, ROW_H, fill=OLIST_BLUE)
        cx = HX
        for i, w in enumerate(cols_w):
            is_emphasis = r[6] and i in (0, 5)
            # 12M 컬럼(i=5)이면서 highlight 행이면 SUCCESS, 아니면 OLIST_BLUE
            if is_emphasis and i == 5:
                color = SUCCESS
            elif is_emphasis:
                color = OLIST_BLUE
            else:
                color = NEAR_BLACK if r[6] else GRAY_600
            font = FONT_BOLD if is_emphasis else FONT_REGULAR
            bold = is_emphasis
            add_text(s, cx + 0.15, ry + 0.11, w - 0.2, 0.3, str(r[i]),
                     size=11, font=font, bold=bold, color=color)
            cx += w
        ry += ROW_H + 0.05

    # 하단 인사이트
    add_rect(s, PAD_X, 6.45, 12.21, 0.55, fill=BLUE_TINT)
    add_rect(s, PAD_X, 6.45, 0.05, 0.55, fill=OLIST_BLUE)
    add_text(s, PAD_X + 0.3, 6.6, 12, 0.3,
             "★ 상단 3행 = TF3 책임. 발표 핵심 KR 모두 TF3에 귀속.",
             size=12, font=FONT_BOLD, bold=True, color=NEAR_BLACK)

    add_footer(s, page_num=25)
    return s


# ============================================================
# Slide 26 — Type 14: Action Plan Roadmap (CTA 4단)
# ============================================================
def make_slide_26(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_header(s, active_idx=3)

    add_page_title(s, runs=[("OLIST  ", OLIST_BLUE),
                            ("STARTS TOMORROW", NEAR_BLACK)], w=7.0, size=42)
    add_text(s, RIGHT_EDGE - 5, 1.20, 5, 0.3,
             "4-PHASE ROLLOUT  ·  Week 1 → Year 1",
             size=9, font=FONT_MEDIUM, color=GRAY_600,
             char_spacing=1.4, align="right")
    add_lead(s, "TF별 4단계 액션 — 이번 주에 시작해 1년 내 KR1 4.5% 달성.")

    # 2D Grid
    GX, GY = PAD_X + 1.2, 2.7
    GW = 12.21 - 1.2
    phases = ["▶ 이번 주", "▶ 이번 달", "▶ BF까지", "▶ BF 이후 1년"]
    col_w = (GW - 3 * 0.15) / 4

    # Header
    for i, p in enumerate(phases):
        cx = GX + i * (col_w + 0.15)
        add_text(s, cx, GY, col_w, 0.3, p,
                 size=10, font=FONT_BOLD, bold=True,
                 color=OLIST_BLUE, char_spacing=1.0)
    add_h_line(s, GX, GY + 0.4, GW, color=GRAY_400)

    tfs = [
        {"label": "TF1\nDATA", "cells": [
            ("배송 지연\n모니터링", "H"),
            ("Fast-Ship\n셀러 300곳", "H"),
            ("BF 검증\n지연 15% 이하", "H"),
            ("3PL 5개 지역\nLT 8일", "H"),
        ], "highlight": False},
        {"label": "TF2\nSEGMENT", "cells": [
            ("PDP A/B\n테스트", "M"),
            ("Metric\nDictionary v1", "H"),
            ("번들 UI\n배포", "H"),
            ("셀러 등급제\n정착", "M"),
        ], "highlight": False},
        {"label": "★ TF3\nTRIGGER", "cells": [
            ("D+7 시나리오\n승인", "H"),
            ("D+7·D+30\nCRM 런칭", "H"),
            ("BF 1.5~2.0%\n검증", "H"),
            ("KR1 4.5%\n달성 ★", "H"),
        ], "highlight": True},
    ]
    ROW_H = 1.0
    ry = GY + 0.55
    for tf in tfs:
        # 좌측 라벨
        label_color = OLIST_BLUE if tf["highlight"] else NEAR_BLACK
        add_text(s, PAD_X, ry + 0.15, 1.15, 0.7, tf["label"],
                 size=11, font=FONT_BOLD, bold=True,
                 color=label_color, line_spacing=1.2)
        # 4 셀
        for i, (text, prio) in enumerate(tf["cells"]):
            cx = GX + i * (col_w + 0.15)
            bg = BLUE_TINT if tf["highlight"] else GRAY_200
            add_rect(s, cx, ry, col_w, ROW_H - 0.1, fill=bg)
            if tf["highlight"]:
                add_rect(s, cx, ry, 0.04, ROW_H - 0.1, fill=OLIST_BLUE)
            add_text(s, cx + 0.15, ry + 0.15, col_w - 0.55, 0.7, text,
                     size=10, font=FONT_BOLD, bold=True,
                     color=NEAR_BLACK, line_spacing=1.3)
            # 우상단 priority badge
            badge_color = OLIST_BLUE if prio == "H" else GRAY_600
            add_rect(s, cx + col_w - 0.32, ry + 0.1, 0.22, 0.22, fill=badge_color)
            add_text(s, cx + col_w - 0.32, ry + 0.13, 0.22, 0.2, prio,
                     size=9, font=FONT_BOLD, bold=True,
                     color=WHITE, align="center")
        ry += ROW_H + 0.05

    # 하단 멘트
    add_text(s, PAD_X, 6.65, 12.21, 0.4,
             "재구매 구조는 선언이 아니라, 이 박스들의 실행에서 만들어집니다.",
             size=14, font=FONT_BOLD, bold=True, color=OLIST_BLUE, align="center")

    add_footer(s, page_num=26)
    return s


# ============================================================
# Slide 27 — Type 8: Looker Studio Dashboard (placeholder)
# ============================================================
def make_slide_27(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_header(s, active_idx=3)

    add_page_title(s, runs=[("LIVE  ", NEAR_BLACK),
                            ("DASHBOARD", OLIST_BLUE)], w=8, size=44)
    add_text(s, RIGHT_EDGE - 5, 1.20, 5, 0.3,
             "AARRR + KR 통합  ·  자료 URL/QR",
             size=9, font=FONT_MEDIUM, color=GRAY_600,
             char_spacing=1.4, align="right")
    add_lead(s, "발표 후 자유롭게 열람 가능 — Looker Studio 대시보드로 실시간 추적.")

    # 좌측 대시보드 — 정교한 mockup (KPI 4 카드 + 차트 placeholder)
    DX, DY, DW, DH = PAD_X, 2.7, 9, 4.0
    add_rect(s, DX, DY, DW, DH, fill=GRAY_200)
    # 대시보드 상단 헤더 띠
    add_rect(s, DX, DY, DW, 0.4, fill=NEAR_BLACK)
    add_text(s, DX + 0.2, DY + 0.08, 4, 0.25,
             "Olist · AARRR Dashboard",
             size=10, font=FONT_BOLD, bold=True, color=WHITE)
    add_text(s, DX + DW - 1.5, DY + 0.08, 1.3, 0.25,
             "2018.08 · 실시간",
             size=8, font=FONT_REGULAR, color=GRAY_400, align="right")
    # KPI 4 카드 (2x2)
    KX, KY = DX + 0.2, DY + 0.6
    KW = (DW - 0.6) / 2
    KH = (DH - 0.7) / 2
    kpis = [
        ("KR1 재구매율", "3.0%", "→ 4.5%", OLIST_BLUE),
        ("BF Cohort 30일", "0.56%", "→ 1.5%", NEAR_BLACK),
        ("KR3 재구매 비중", "6.13%", "→ 9.0%", NEAR_BLACK),
        ("KR4 첫 구매 지연", "8.16%", "→ 6.0%", NEAR_BLACK),
    ]
    for i, (label, current, target, color) in enumerate(kpis):
        mx = KX + (i % 2) * (KW + 0.2)
        my = KY + (i // 2) * (KH + 0.2)
        add_rect(s, mx, my, KW, KH, fill=WHITE)
        if color == OLIST_BLUE:
            add_rect(s, mx, my, 0.05, KH, fill=OLIST_BLUE)
        add_text(s, mx + 0.2, my + 0.15, KW - 0.4, 0.25, label,
                 size=8, font=FONT_MEDIUM, color=GRAY_600, char_spacing=1.0)
        add_text(s, mx + 0.2, my + 0.45, KW - 0.4, 0.6, current,
                 size=28, font=FONT_BLACK, bold=True, color=NEAR_BLACK)
        add_text(s, mx + 0.2, my + 1.05, KW - 0.4, 0.3, target,
                 size=11, font=FONT_BOLD, bold=True, color=color)

    # 우측 QR 카드
    QX, QY, QW, QH = 9.7, 2.7, 2.83, 4.0
    add_rect(s, QX, QY, QW, QH, fill=BLUE_TINT)
    add_rect(s, QX, QY, 0.05, QH, fill=OLIST_BLUE)
    add_text(s, QX + 0.3, QY + 0.3, QW - 0.6, 0.3,
             "ACCESS", size=10, font=FONT_MEDIUM,
             color=OLIST_BLUE, char_spacing=1.4)
    # QR placeholder
    qr_size = 1.9
    qr_x = QX + (QW - qr_size) / 2
    qr_y = QY + 0.75
    add_rect(s, qr_x, qr_y, qr_size, qr_size, fill=WHITE, line=GRAY_400)
    add_text(s, qr_x, qr_y + qr_size / 2 - 0.15, qr_size, 0.3,
             "[QR]", size=14, font=FONT_BOLD, bold=True,
             color=GRAY_400, align="center")
    add_text(s, QX + 0.3, QY + 2.85, QW - 0.6, 0.3,
             "Scan to access",
             size=11, font=FONT_BOLD, bold=True,
             color=NEAR_BLACK, align="center")
    add_text(s, QX + 0.3, QY + 3.2, QW - 0.6, 0.5,
             "lookerstudio.\ngoogle.com/...",
             size=9, font=FONT_REGULAR, color=GRAY_600,
             align="center", line_spacing=1.3)

    # 하단
    add_text(s, PAD_X, 6.85, 12.21, 0.25,
             "발표 중에는 패스 — 발표 후 자유롭게 열람 가능합니다.",
             size=11, font=FONT_REGULAR, color=GRAY_600, align="center")

    add_footer(s, page_num=27)
    return s


# ============================================================
# Slide 28 — Type 4: Further Reading (목록형)
# ============================================================
def make_slide_28(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_header(s, active_idx=3)

    add_page_title(s, runs=[("FURTHER  ", NEAR_BLACK),
                            ("READING", OLIST_BLUE)], size=44)
    add_lead(s, "발표에서 다 못 보여드린 분석들 — Q&A 대응용.")

    # 좌단: 메인 분석
    LX, LY = PAD_X, 2.7
    LW = 5.85
    add_text(s, LX, LY, LW, 0.3, "메인 분석",
             size=10, font=FONT_MEDIUM, color=GRAY_600, char_spacing=1.4)
    add_h_line(s, LX, LY + 0.35, LW)
    main_items = [
        ("01", "지역×카테고리 교차분석", "RFM × 카테고리 매트릭스"),
        ("02", "리뷰 감성분석", "1점 리뷰 키워드 클러스터"),
        ("03", "ARR 상세 지표 테이블", "코호트별 LTV·CAC"),
        ("04", "H2 가설 검증", "5점 편향 통계적 유의성"),
        ("05", "코호트 생존 곡선", "월 코호트별 잔존율"),
        ("06", "TF4 (셀러 통합) · TF5 (측정 통일)", "확장 TF 후보"),
    ]
    iy = LY + 0.55
    for num, title, sub in main_items:
        add_text(s, LX + 0.1, iy, 0.5, 0.3, num,
                 size=11, font=FONT_BLACK, bold=True, color=OLIST_BLUE)
        add_text(s, LX + 0.7, iy, LW - 0.8, 0.3, title,
                 size=12, font=FONT_BOLD, bold=True, color=NEAR_BLACK)
        add_text(s, LX + 0.7, iy + 0.3, LW - 0.8, 0.25, sub,
                 size=10, font=FONT_REGULAR, color=GRAY_600)
        iy += 0.6

    # 우단: 부록
    RX = LX + LW + 0.5
    RW = LW
    add_text(s, RX, LY, RW, 0.3, "부록 (Q&A 대비)",
             size=10, font=FONT_MEDIUM, color=GRAY_600, char_spacing=1.4)
    add_h_line(s, RX, LY + 0.35, RW)
    appx_items = [
        ("A-1", "AI 파이프라인 + 팀원 R&R", "BigQuery → Claude → Looker → Pencil"),
        ("A-2", "회고", "어려웠던 점 / 배운 점 / 다음에 더 잘할 것"),
        ("A-4", "TF 충돌 사례 카드", "BF cohort 9배 차이 등"),
    ]
    iy = LY + 0.55
    for num, title, sub in appx_items:
        add_text(s, RX + 0.1, iy, 0.5, 0.3, num,
                 size=11, font=FONT_BLACK, bold=True, color=OLIST_BLUE)
        add_text(s, RX + 0.7, iy, RW - 0.8, 0.3, title,
                 size=12, font=FONT_BOLD, bold=True, color=NEAR_BLACK)
        add_text(s, RX + 0.7, iy + 0.3, RW - 0.8, 0.25, sub,
                 size=10, font=FONT_REGULAR, color=GRAY_600)
        iy += 0.6

    # 우단 하단 카드
    iy += 0.25
    add_rect(s, RX, iy, RW, 1.4, fill=BLUE_TINT)
    add_rect(s, RX, iy, 0.05, 1.4, fill=OLIST_BLUE)
    add_text(s, RX + 0.3, iy + 0.2, RW - 0.6, 0.3,
             "Q&A 대비 추가 자료",
             size=10, font=FONT_MEDIUM, color=OLIST_BLUE, char_spacing=1.4)
    add_text(s, RX + 0.3, iy + 0.55, RW - 0.6, 0.7,
             "위 자료는 모두 슬27 Looker 대시보드와\n연동되어 실시간 추적 가능합니다.",
             size=11, font=FONT_REGULAR, color=NEAR_BLACK, line_spacing=1.5)

    add_footer(s, page_num=28)
    return s


# ============================================================
# Slide 29 (NEW) — Appendix: TF3 CRM 9 Channels
# ============================================================
def make_slide_appendix_crm(prs, page_num=29):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_header(s, active_idx=3)  # Outcome

    add_page_title(s, runs=[("APPENDIX  ", NEAR_BLACK),
                            ("TF3 CRM 9 CHANNELS", OLIST_BLUE)],
                   w=11, size=36)
    add_lead(s, "Q&A 대응용 — TF3 CRM 9채널 시안 풀(D+7 / D+30 / D+32).")

    # 그리드: 좌측 라벨(1.0") + 4열 그리드
    GX, GY = PAD_X + 1.0, 2.6
    GW = 12.21 - 1.0
    cols = 4
    col_w = (GW - 3 * 0.12) / cols
    row_h = 1.4
    row_gap = 0.15

    rows = [
        {"period": "D+7", "items": [
            ("Push A", "발표자료/TF3-재구매/assets/tf3_D07_push_A.png"),
            ("Push B", "발표자료/TF3-재구매/assets/tf3_D07_push_B.png"),
            ("Email",  "발표자료/TF3-재구매/assets/tf3_D07_email.png"),
            ("Kakao",  "발표자료/TF3-재구매/assets/tf3_D07_kakao.png"),
        ]},
        {"period": "D+30", "items": [
            ("Push A", "발표자료/TF3-재구매/assets/tf3_D30_push_A.png"),
            ("Push B", "발표자료/TF3-재구매/assets/tf3_D30_push_B.png"),
            ("Email",  "발표자료/TF3-재구매/assets/tf3_D30_email.png"),
            ("Kakao ★","발표자료/TF3-재구매/assets/tf3_D30_kakao.png"),
        ]},
        {"period": "D+32", "items": [
            ("Kakao ★", "발표자료/TF3-재구매/assets/tf3_D32_kakao.png"),
        ]},
    ]

    ry = GY
    for row in rows:
        # 좌측 시점 라벨
        add_text(s, PAD_X, ry + 0.3, 1.0, 0.4, row["period"],
                 size=18, font=FONT_BLACK, bold=True, color=OLIST_BLUE)
        for i, (label, img_path) in enumerate(row["items"]):
            cx = GX + i * (col_w + 0.12)
            # 채널 라벨
            add_text(s, cx, ry, col_w, 0.2, label,
                     size=8, font=FONT_BOLD, bold=True,
                     color=NEAR_BLACK, char_spacing=0.6, align="center")
            # 이미지
            add_picture_safe(s, img_path, cx, ry + 0.2, col_w, row_h - 0.2)
        ry += row_h + row_gap

    # 하단 노트
    add_rect(s, PAD_X, 6.65, 12.21, 0.45, fill=BLUE_TINT)
    add_rect(s, PAD_X, 6.65, 0.05, 0.45, fill=OLIST_BLUE)
    add_text(s, PAD_X + 0.3, 6.78, 12, 0.25,
             "★ 발표 본편(슬23)에 노출된 best-of 3장. 나머지는 A/B 비교 + 채널 다양화 검증용.",
             size=10, font=FONT_REGULAR, color=NEAR_BLACK)

    add_footer(s, page_num=page_num)
    return s


# ============================================================
# Slide 31 (이전 30) — Type 9 변형: Q&A (다크 풀블리드)
# ============================================================
def make_slide_30(prs, page_num=31):
    s = prs.slides.add_slide(prs.slide_layouts[6])

    bg = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0,
                              Inches(SLIDE_W), Inches(SLIDE_H))
    bg.fill.solid()
    bg.fill.fore_color.rgb = OLIST_BLUE
    bg.line.fill.background()
    bg.shadow.inherit = False

    add_header(s, active_idx=-1, dark=True)

    # Mega "Q & A" — 폰트 180pt + 박스 3.0 (슬9·슬29 패치 패턴 적용)
    add_text(s, 0, 1.8, SLIDE_W, 3.0, "Q & A",
             size=180, font=FONT_BLACK, color=WHITE,
             align="center", char_spacing=0, anchor="middle")

    # Lead
    add_text(s, 0, 4.95, SLIDE_W, 0.4,
             "질문 받겠습니다.",
             size=14, font=FONT_LIGHT,
             color=RGBColor(0xCC, 0xD3, 0xFF), align="center")

    # Bottom
    add_text(s, 0, 5.85, SLIDE_W, 0.3,
             "데2터로말해조  ·  growth@olist.com",
             size=10, font=FONT_MEDIUM,
             color=RGBColor(0xCC, 0xD3, 0xFF), align="center")

    add_footer(s, page_num=page_num, dark=True)
    return s


# ============================================================
# Main — 30장 전체
# ============================================================
def main():
    out_dir = Path(__file__).resolve().parent.parent / "exports"
    out_dir.mkdir(exist_ok=True)
    out_path = out_dir / "olist_full_v2.pptx"

    prs = Presentation()
    prs.slide_width = Inches(SLIDE_W)
    prs.slide_height = Inches(SLIDE_H)

    slides_to_make = [
        (1,  make_slide_1,  "Cover"),
        (2,  make_slide_2,  "Email Card"),
        (3,  make_slide_3,  "Team 4-Quadrant"),
        (4,  make_slide_4,  "Executive Summary (SCR)"),
        (5,  make_slide_5,  "3-Axis Signal"),
        (6,  make_slide_6,  "RFM Segmentation"),
        (7,  make_slide_7,  "Hypothesis Tree"),
        (8,  make_slide_8,  "Low Rating Distribution"),
        (9,  make_slide_9,  "Quote/Stat (Peak)"),
        (10, make_slide_10, "Findings"),
        (11, make_slide_11, "Why TF"),
        (12, make_slide_12, "5x3 Matrix"),
        (13, make_slide_13, "Dual Track"),
        (14, make_slide_14, "TF1 KPI"),
        (15, make_slide_15, "TF1 B/A"),
        (16, make_slide_16, "TF1 Outputs"),
        (17, make_slide_17, "TF2 KPI"),
        (18, make_slide_18, "TF2 B/A"),
        (19, make_slide_19, "TF2 Outputs"),
        (20, make_slide_20, "TF3 KPI"),
        (21, make_slide_21, "TF3 B/A"),
        (22, make_slide_22, "Lifecycle Journey Map"),
        (23, make_slide_23, "TF3 Outputs (NEW)"),
        (24, make_slide_24, "Conclusion (NEW)"),
        (25, make_slide_25, "Expected Impact (NEW)"),
        (26, make_slide_26, "Action Roadmap (NEW)"),
        (27, make_slide_27, "Looker Dashboard (NEW)"),
        (28, make_slide_28, "Further Reading"),
        (29, make_slide_appendix_crm, "Appendix: TF3 CRM 9 Channels (NEW)"),
        (30, make_slide_29, "THANK YOU"),
        (31, make_slide_30, "Q&A"),
    ]
    for num, fn, label in slides_to_make:
        # 슬29(부록), 30(THANK YOU), 31(Q&A)는 page_num 인자 받음
        if num in (29, 30, 31):
            print(f"  Slide {num:2d} - {label}")
            fn(prs, page_num=num)
        else:
            print(f"  Slide {num:2d} - {label}")
            fn(prs)

    prs.save(out_path)
    print(f"\nDone -> {out_path}")
    print(f"  size: {out_path.stat().st_size:,} bytes")


if __name__ == "__main__":
    main()
