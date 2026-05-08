"""PPTX 무결성 + 디자인 토큰 검증 (디자인시스템 v1.2 §11.7)"""

import zipfile
import re
import sys
from pathlib import Path
from pptx import Presentation

def validate(p: Path):
    print(f"=== Validating {p.name} ===\n")

    # 1. python-pptx 로드
    try:
        prs = Presentation(p)
        print(f"[OK] Loadable: {len(prs.slides)} slides")
        print(f"     Size: {prs.slide_width.inches}\" x {prs.slide_height.inches}\"")
    except Exception as e:
        print(f"[FAIL] Load error: {e}")
        return False

    # 2. LINE 도형 cx=0 / cy=0 검사
    issues = []
    with zipfile.ZipFile(p) as z:
        slide_files = sorted([n for n in z.namelist() if n.startswith("ppt/slides/slide") and n.endswith(".xml")])
        print(f"     XML slides: {len(slide_files)}")
        for sf in slide_files:
            content = z.read(sf).decode("utf-8")
            for cx, cy in re.findall(r'<a:ext cx="(\d+)" cy="(\d+)"', content):
                if int(cx) == 0 or int(cy) == 0:
                    issues.append(f"{sf}: cx={cx} cy={cy}")

    if issues:
        print(f"[FAIL] LINE shapes with cx=0 or cy=0:")
        for i in issues:
            print(f"  {i}")
        return False
    else:
        print(f"[OK] No LINE shape with cx=0 or cy=0")

    # 3. 폰트 검사 — Pretendard 외 사용 여부
    fonts = set()
    with zipfile.ZipFile(p) as z:
        for sf in slide_files:
            content = z.read(sf).decode("utf-8")
            for f in re.findall(r'typeface="([^"]+)"', content):
                fonts.add(f)
    non_pretendard = [f for f in fonts if "Pretendard" not in f and f not in ("+mj-lt", "+mn-lt", "+mj-ea", "+mn-ea")]
    print(f"     Fonts used: {sorted(fonts)}")
    if non_pretendard:
        print(f"[WARN] Non-Pretendard fonts: {non_pretendard}")
    else:
        print(f"[OK] Only Pretendard family used")

    # 4. 슬라이드별 도형 수
    print(f"\n     Shape counts:")
    for i, slide in enumerate(prs.slides, 1):
        print(f"       Slide {i}: {len(slide.shapes)} shapes")

    return True

if __name__ == "__main__":
    target = Path(sys.argv[1]) if len(sys.argv) > 1 else \
             Path(__file__).resolve().parent.parent / "exports" / "olist_demo_v1.pptx"
    ok = validate(target)
    sys.exit(0 if ok else 1)
